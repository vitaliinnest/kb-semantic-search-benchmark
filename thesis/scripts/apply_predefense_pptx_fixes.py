"""
Apply pre-defense slide text fixes:
- Slide 2: replace «Вибір оптимальної моделі» (task 6) with concrete name
- Slide 4: bind «обґрунтовано визначити модель» to MCDA method;
            replace vague criterion «якість семантичного пошуку»
- Slide 15: correct the «3 доменах» factual claim about BGE-M3 dominance
- Slide 16: replace informal «стеля якості» with neutral phrasing
- Slide 17: remove «РОЗДІЛ 11 ·» label + tighten conclusion wording +
            quantify «прийнятна якість» of E5-base
- Slide 18: remove «РОЗДІЛ 10 ·» label

Walks all shapes/paragraphs/runs and replaces within whichever run holds
the substring. Falls back to merging runs if the target text spans
multiple runs.
"""
import sys
from pathlib import Path
from pptx import Presentation

PPTX = Path(__file__).parent.parent / "2026_М_ПІ_ІПЗм-24-1_Нестеренко_В_В.pptx"

EDITS: dict[int, list[tuple[str, str]]] = {
    2: [
        ("Вибір оптимальної моделі", "Багатокритеріальний вибір моделі"),
        (
            "визначення доцільності застосування сучасних моделей векторних ембеддінгів для підвищення ефективності семантичного пошуку текстових документів у корпоративних базах знань шляхом порівняльного експериментального дослідження.",
            "визначити найбільш придатну embedding-модель для україномовного семантичного пошуку у корпоративних базах знань на основі експериментального порівняння та багатокритеріального аналізу.",
        ),
    ],
    4: [
        (
            "необхідно обґрунтовано визначити модель, що забезпечує оптимальний баланс",
            "необхідно обґрунтовано обрати модель методом багатокритеріального аналізу за балансом",
        ),
        ("якість семантичного пошуку", "якість ранжування (nDCG@10)"),
    ],
    15: [
        (
            "лідер у 3 доменах за nDCG, MRR, Recall",
            "лідер за nDCG, MRR, Recall на технічному і медичному доменах",
        ),
    ],
    16: [
        ("стеля якості ~0.49 nDCG", "найвища nDCG ~0.49"),
    ],
    17: [
        ("РОЗДІЛ 11 · ПІДСУМКИ", "ПІДСУМКИ"),
        (
            "Обґрунтовано доцільність застосування embedding-моделей для семантичного пошуку у корпоративних базах знань",
            "Embedding-моделі статистично значуще переважають BM25 для україномовного семантичного пошуку",
        ),
        (
            "Прийнятна якість + ~3× швидше за BGE-M3",
            "Якість на ~10% нижча за BGE-M3, ~3× швидше",
        ),
    ],
    18: [
        ("РОЗДІЛ 10 · ПУБЛІКАЦІЯ РЕЗУЛЬТАТІВ", "АПРОБАЦІЯ РЕЗУЛЬТАТІВ"),
    ],
}

# Same structure but applied to speaker notes text frames instead of slide text.
# Used to keep voice-over in sync with slide edits above.
NOTES_EDITS: dict[int, list[tuple[str, str]]] = {
    15: [
        (
            "BGE-M3 — лідер у трьох доменах за nDCG, MRR, Recall.",
            "BGE-M3 — лідер за nDCG, MRR, Recall на технічному і медичному доменах.",
        ),
    ],
    17: [
        (
            "Підсумовуючи: у роботі обґрунтовано доцільність застосування сучасних embedding-моделей для семантичного пошуку у корпоративних базах знань. Сформульовано практичні рекомендації щодо вибору моделі.",
            "Підсумовуючи: embedding-моделі статистично значуще переважають BM25 для україномовного семантичного пошуку. Сформульовано практичні рекомендації щодо вибору моделі під конкретні умови.",
        ),
        (
            "E5-base — швидка альтернатива з прийнятною якістю та триразовою перевагою за швидкодією.",
            "E5-base — швидка альтернатива: якість на ~10% нижча за BGE-M3, але утричі швидше.",
        ),
    ],
}


def apply_text_edit(slide, old: str, new: str) -> bool:
    """Replace `old` with `new` in any text run on `slide`."""
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        for para in shape.text_frame.paragraphs:
            # Cleanest path: substring sits in a single run
            for run in para.runs:
                if old in run.text:
                    run.text = run.text.replace(old, new)
                    return True
            # Fallback: substring spans multiple runs — merge into first
            combined = "".join(run.text for run in para.runs)
            if old in combined:
                new_combined = combined.replace(old, new)
                if para.runs:
                    para.runs[0].text = new_combined
                    for r in para.runs[1:]:
                        r.text = ""
                return True
    return False


def apply_notes_edit(slide, old: str, new: str) -> bool:
    """Replace `old` with `new` in the slide's speaker-notes text frame."""
    if not slide.has_notes_slide:
        return False
    tf = slide.notes_slide.notes_text_frame
    for para in tf.paragraphs:
        for run in para.runs:
            if old in run.text:
                run.text = run.text.replace(old, new)
                return True
        combined = "".join(run.text for run in para.runs)
        if old in combined:
            new_combined = combined.replace(old, new)
            if para.runs:
                para.runs[0].text = new_combined
                for r in para.runs[1:]:
                    r.text = ""
            return True
    return False


def main() -> None:
    sys.stdout.reconfigure(encoding="utf-8")
    prs = Presentation(str(PPTX))
    for slide_num, edits in EDITS.items():
        slide = prs.slides[slide_num - 1]
        print(f"=== Slide {slide_num} (text) ===")
        for old, new in edits:
            ok = apply_text_edit(slide, old, new)
            status = "OK  " if ok else "MISS"
            print(f"  [{status}] {old[:70]}")
    for slide_num, edits in NOTES_EDITS.items():
        slide = prs.slides[slide_num - 1]
        print(f"=== Slide {slide_num} (notes) ===")
        for old, new in edits:
            ok = apply_notes_edit(slide, old, new)
            status = "OK  " if ok else "MISS"
            print(f"  [{status}] {old[:70]}")
    prs.save(str(PPTX))
    print(f"\nSaved: {PPTX}")


if __name__ == "__main__":
    main()
