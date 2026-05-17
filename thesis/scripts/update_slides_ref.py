"""
Update thesis/SLIDES_REF.md after slide 13 deletion:
- remove the old slide 13 section
- renumber slides 14..18 → 13..17
- refresh all time intervals from the actual .pptx via pptx module
- update the document header «18 слайдах» → «17 слайдах»
"""
import sys
import re
from pathlib import Path
from pptx import Presentation

REF = Path(__file__).parent.parent / "SLIDES_REF.md"
PPTX = Path(__file__).parent.parent / "2026_М_ПІ_ІПЗм-24-1_Нестеренко_В_В.pptx"


def fetch_intervals_from_pptx() -> list[str]:
    """Return list of «M:SS – M:SS» strings, one per slide, from current pptx state."""
    prs = Presentation(str(PPTX))
    intervals: list[str] = []
    for slide in prs.slides:
        tf = slide.notes_slide.notes_text_frame if slide.has_notes_slide else None
        if tf is None:
            intervals.append("")
            continue
        for para in tf.paragraphs:
            if para.text.startswith("⏱"):
                # Format: «⏱ M:SS – M:SS  ·  слайд N/17»
                m = re.search(r"(\d+:\d{2})\s*[–-]\s*(\d+:\d{2})", para.text)
                if m:
                    intervals.append(f"{m.group(1)} – {m.group(2)}")
                else:
                    intervals.append("")
                break
        else:
            intervals.append("")
    return intervals


def main() -> None:
    sys.stdout.reconfigure(encoding="utf-8")
    intervals = fetch_intervals_from_pptx()
    assert len(intervals) == 17, f"Expected 17 slides, got {len(intervals)}"
    print(f"Fetched {len(intervals)} intervals from pptx")

    content = REF.read_text(encoding="utf-8")

    # 1. Update header
    content = content.replace(
        "довідник по 18 слайдах",
        "довідник по 17 слайдах",
    )

    # 2. Locate and delete the old slide 13 section («## Слайд 13 — ...»)
    #    The section starts at «## Слайд 13` and ends right before the next «## Слайд 14»
    slide13_start = content.find("## Слайд 13")
    slide14_start = content.find("## Слайд 14")
    assert slide13_start > 0 and slide14_start > slide13_start, "Could not locate slide 13/14 boundaries"

    # The «---\n\n» separator before slide 14 belongs to the boundary — remove it as well
    # Actually keep the separator structure: section ends at the «---\n\n## Слайд 14»
    # We'll remove from slide13_start up to but not including slide14_start
    content = content[:slide13_start] + content[slide14_start:]
    print("Removed old slide 13 section")

    # 3. Renumber «## Слайд N — ...» for N=14..18 → 13..17 (work in reverse to avoid collisions)
    for old_n in range(18, 13, -1):
        new_n = old_n - 1
        old_marker = f"## Слайд {old_n} —"
        new_marker = f"## Слайд {new_n} —"
        before = content.count(old_marker)
        content = content.replace(old_marker, new_marker)
        after = content.count(new_marker)
        print(f"  renumbered slide {old_n} → {new_n} ({before} match(es))")

    # 4. Refresh time intervals on all 17 slide headings.
    #    Pattern: «## Слайд N — Title `X:XX – Y:YY`»
    def replace_interval(match: re.Match) -> str:
        n = int(match.group(1))
        title = match.group(2)
        new_interval = intervals[n - 1]
        return f"## Слайд {n} — {title} `{new_interval}`"

    content = re.sub(
        r"## Слайд (\d+) — (.+?) `(\d+:\d{2}\s*[–-]\s*\d+:\d{2})`",
        replace_interval,
        content,
    )
    print("Refreshed all 17 time intervals")

    REF.write_text(content, encoding="utf-8")
    print(f"\nSaved: {REF}")


if __name__ == "__main__":
    main()
