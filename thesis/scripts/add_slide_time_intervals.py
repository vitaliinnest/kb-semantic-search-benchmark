"""
Replace the first line of each slide's speaker notes
    «⏱ ~XX сек  ·  слайд N/18»
with a cumulative-time interval
    «⏱ M:SS – M:SS  ·  слайд N/18»
based on actual word count (assuming 2.3 words/sec reading rate).

Total fits within the 10-minute presentation budget.
"""
import re
import sys
from pathlib import Path
from pptx import Presentation

PPTX = Path(__file__).parent.parent / "2026_М_ПІ_ІПЗм-24-1_Нестеренко_В_В.pptx"
WPS = 2.3  # words per second (academic reading rate)


def fmt(seconds: float) -> str:
    s = round(seconds)
    return f"{s // 60}:{s % 60:02d}"


def main() -> None:
    sys.stdout.reconfigure(encoding="utf-8")
    prs = Presentation(str(PPTX))
    n_slides = len(prs.slides)

    # First pass: compute per-slide word counts (excluding the timing first line)
    durations: list[float] = []
    note_texts: list[str] = []
    for slide in prs.slides:
        notes = slide.notes_slide.notes_text_frame.text if slide.has_notes_slide else ""
        note_texts.append(notes)
        body_lines = [ln for ln in notes.splitlines() if not ln.startswith("⏱")]
        body = "\n".join(body_lines).strip()
        words = len(body.split())
        durations.append(words / WPS)

    # Cumulative intervals
    cum = 0.0
    intervals: list[tuple[str, str]] = []
    for d in durations:
        start = cum
        cum += d
        intervals.append((fmt(start), fmt(cum)))

    print(f"Total: {fmt(cum)} (target: under 10:00)")
    print()
    for i, (start, end) in enumerate(intervals, 1):
        print(f"  слайд {i:2}/{n_slides}: {start} – {end}")
    print()

    # Second pass: rewrite the first line
    for i, (slide, (start, end)) in enumerate(zip(prs.slides, intervals), 1):
        if not slide.has_notes_slide:
            continue
        tf = slide.notes_slide.notes_text_frame
        new_first_line = f"⏱ {start} – {end}  ·  слайд {i}/{n_slides}"

        # Build new text: replace the line that starts with ⏱, keep rest
        old_text = note_texts[i - 1]
        old_lines = old_text.splitlines()
        new_lines: list[str] = []
        replaced = False
        for ln in old_lines:
            if ln.startswith("⏱") and not replaced:
                new_lines.append(new_first_line)
                replaced = True
            else:
                new_lines.append(ln)
        if not replaced:
            new_lines.insert(0, new_first_line)
            new_lines.insert(1, "")

        # Write back: clear existing paragraphs, recreate
        tf.clear()
        for j, ln in enumerate(new_lines):
            if j == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = ln

    prs.save(str(PPTX))
    print(f"Saved: {PPTX}")


if __name__ == "__main__":
    main()
