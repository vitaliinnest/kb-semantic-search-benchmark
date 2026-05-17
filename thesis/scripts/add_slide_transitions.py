"""
Add short bridging phrases at the start of each slide's notes so the
voice-over flows from one slide to the next instead of starting cold.

The closing «Дякую за увагу» line is intentionally kept out — it's a
live-defense thing and the student prefers to say it ad-lib.

Run after this: add_slide_time_intervals.py to refresh ⏱ markers.
"""
import sys
from pathlib import Path
from pptx import Presentation

PPTX = Path(__file__).parent.parent / "2026_М_ПІ_ІПЗм-24-1_Нестеренко_В_В.pptx"

# slide_num → prefix to insert at the start of the first body paragraph
TRANSITIONS: dict[int, str] = {
    2:  "Розпочнемо з актуальності. ",
    3:  "Які підходи вже досліджено в літературі? ",
    4:  "Виявлена прогалина веде до постановки задачі. ",
    5:  "Розглянемо, які саме моделі дослідимо. ",
    6:  "Як порівнюватимемо ці п'ять моделей? ",
    8:  "Окремих метрик недостатньо — потрібно об'єднати їх в одне число. ",
    9:  "Як це реалізовано архітектурно? ",
    10: "Технологічний стек, на якому це побудовано. ",
    11: "На якому датасеті випробували методологію? ",
    13: "Що дав запуск benchmark на 300 запитах? ",
    14: "За окремими доменами картина така. ",
    15: "Переходимо до інтегральної оцінки. ",
    16: "Той самий аналіз — у вигляді графіка. ",
    18: "Окремо — про апробацію результатів роботи. ",
}

# Slide 17 still has «Дякую за увагу» from prior edits — drop it if found
SLIDE_17_REMOVE = "Дякую за увагу! Готовий відповісти на запитання."


def first_body_para_idx(tf) -> int:
    """Index of the first paragraph that isn't the ⏱ timing line."""
    for i, p in enumerate(tf.paragraphs):
        if not p.text.startswith("⏱") and p.text.strip():
            return i
    return -1


def prepend_to_para(para, text: str) -> None:
    """Prepend text to a paragraph, preserving formatting of first run."""
    if para.runs:
        para.runs[0].text = text + para.runs[0].text
    else:
        para.text = text + para.text


def remove_paragraph_containing(tf, needle: str) -> bool:
    """Remove the first paragraph whose text contains needle."""
    for p in tf.paragraphs:
        if needle in p.text:
            p._p.getparent().remove(p._p)
            return True
    return False


def main() -> None:
    sys.stdout.reconfigure(encoding="utf-8")
    prs = Presentation(str(PPTX))

    for slide_num, prefix in TRANSITIONS.items():
        slide = prs.slides[slide_num - 1]
        tf = slide.notes_slide.notes_text_frame
        idx = first_body_para_idx(tf)
        if idx < 0:
            print(f"[MISS] slide {slide_num}: no body paragraph")
            continue
        # Idempotent: skip if the transition is already there
        if tf.paragraphs[idx].text.startswith(prefix.strip()):
            print(f"[SKIP] slide {slide_num}: already has transition")
            continue
        prepend_to_para(tf.paragraphs[idx], prefix)
        print(f"[OK  ] slide {slide_num}: + «{prefix.strip()}»")

    # Slide 17: drop the thanks line if it's still there from prior edits
    slide17 = prs.slides[16]
    tf17 = slide17.notes_slide.notes_text_frame
    if remove_paragraph_containing(tf17, SLIDE_17_REMOVE):
        print(f"[OK  ] slide 17: removed thanks line")
    else:
        print(f"[SKIP] slide 17: thanks line already gone")

    prs.save(str(PPTX))
    print(f"\nSaved: {PPTX}")


if __name__ == "__main__":
    main()
