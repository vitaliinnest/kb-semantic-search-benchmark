"""
Pre-defense fixes for Nesterenko_Presentation.pptx.

Surgical edits — does not regenerate slides from scratch.

Changes:
  Slide 7  — DCG formula: linear -> exponential (match thesis)
  Slide 14 — Latency note: real numbers from results JSON (not 203-2960 etc.)
  Slide 14 — Chart axis tick labels: shrink to fit
  Slide 15 — Findings: "~225 vs ~2960", "~1500 GPU" -> real numbers + Qwen3 Legal-win
  Slide 16 — Scatter chart data: replace with actual avg_latency_ms from JSON
  Slide 16 — Takeaway: "Оптимум: E5-base ~225" -> "Альтернатива: E5-base ~70"
  Slide 18 — Recommendations: "13× faster" -> "~3×", "Qwen3 GPU only" -> domain-specific
"""
import sys
import pathlib

sys.stdout.reconfigure(encoding="utf-8")

from pptx import Presentation
from pptx.util import Pt, Inches
from pptx.dml.color import RGBColor
from pptx.chart.data import XyChartData

ROOT = pathlib.Path("D:/repos/kb-semantic-search-benchmark")
PPTX = ROOT / "thesis" / "Nesterenko_Presentation.pptx"


def find_run(slide, substring):
    """Return (shape, paragraph, run) for first run containing substring, else None."""
    for sh in slide.shapes:
        if not sh.has_text_frame:
            continue
        for p in sh.text_frame.paragraphs:
            for r in p.runs:
                if substring in r.text:
                    return sh, p, r
    return None


def replace_run_text(slide, needle, new_text):
    hit = find_run(slide, needle)
    if hit is None:
        print(f"  ! not found: {needle!r}")
        return False
    _, _, r = hit
    r.text = new_text
    return True


def main():
    prs = Presentation(str(PPTX))
    s = prs.slides

    # ── Slide 7 — DCG formula ──────────────────────────────────────────
    print("Slide 7: DCG formula -> exponential")
    replace_run_text(
        s[6],
        "де DCG@k = Σᵢ relᵢ / log₂(i+1)",
        "де DCG@k = Σᵢ (2^relᵢ − 1) / log₂(i+1),\nIDCG@k — DCG ідеального ранжування",
    )

    # ── Slide 14 — latency note ────────────────────────────────────────
    print("Slide 14: latency note + chart axis labels")
    replace_run_text(
        s[13],
        "BGE-M3 ≈ 203–2960",
        "BM25 ≈ 0.4–3.1   ◇   E5-base ≈ 63–83   ◇   nomic ≈ 125–168   "
        "◇   BGE-M3 ≈ 208–255   ◇   Qwen3 ≈ 436–463",
    )
    # Shrink chart axis tick labels (Тех./Юр./Мед. were huge)
    for sh in s[13].shapes:
        if sh.has_chart:
            chart = sh.chart
            try:
                chart.category_axis.tick_labels.font.size = Pt(9)
                chart.value_axis.tick_labels.font.size = Pt(8)
            except Exception as e:
                print(f"  ! axis label resize failed: {e}")

    # ── Slide 15 — МКВ findings text ───────────────────────────────────
    print("Slide 15: МКВ findings — real latency + Qwen3 Legal-win")
    replace_run_text(
        s[14],
        "швидка альтернатива (~225 мс vs ~2960 мс у BGE-M3)",
        "швидка альтернатива (~70 мс vs ~228 мс у BGE-M3)",
    )
    replace_run_text(
        s[14],
        "конкурентна якість, але ~1500 мс на CPU — потрібен GPU",
        "лідер на юридичному домені (nDCG@10=0.320), але ~2× повільніший за BGE-M3",
    )

    # ── Slide 15 — bar chart axis labels ──────────────────────────────
    print("Slide 15: bar chart axis labels — shrink")
    for sh in s[14].shapes:
        if sh.has_chart:
            try:
                sh.chart.category_axis.tick_labels.font.size = Pt(11)
                sh.chart.value_axis.tick_labels.font.size = Pt(9)
            except Exception as e:
                print(f"  ! axis label resize failed: {e}")

    # ── Slide 16 — scatter chart data ──────────────────────────────────
    print("Slide 16: scatter chart — real ms values")
    chart_replaced = False
    for sh in s[15].shapes:
        if sh.has_chart and sh.chart.chart_type.__str__().startswith("XY"):
            chart = sh.chart
            cd = XyChartData()
            bge = cd.add_series("BGE-M3")
            for ms, ndcg in [(222.5, 0.6722), (208.2, 0.3065), (254.9, 0.4339)]:
                bge.add_data_point(ms, ndcg)
            e5 = cd.add_series("E5-base")
            for ms, ndcg in [(70.9, 0.6121), (62.5, 0.2567), (83.3, 0.3909)]:
                e5.add_data_point(ms, ndcg)
            qwen = cd.add_series("Qwen3")
            for ms, ndcg in [(463.4, 0.6325), (436.7, 0.3199), (436.2, 0.3629)]:
                qwen.add_data_point(ms, ndcg)
            bm25 = cd.add_series("BM25")
            for ms, ndcg in [(0.4, 0.4861), (3.1, 0.1875), (0.6, 0.3222)]:
                bm25.add_data_point(ms, ndcg)
            nom = cd.add_series("nomic")
            for ms, ndcg in [(124.8, 0.3765), (131.5, 0.0951), (167.9, 0.1668)]:
                nom.add_data_point(ms, ndcg)
            chart.replace_data(cd)
            try:
                chart.category_axis.tick_labels.font.size = Pt(9)
                chart.value_axis.tick_labels.font.size = Pt(9)
            except Exception:
                pass
            chart_replaced = True
            break
    if not chart_replaced:
        print("  ! XY scatter chart not found on slide 16")

    # ── Slide 16 — takeaways ───────────────────────────────────────────
    print("Slide 16: takeaways — Оптимум -> Альтернатива, real ms")
    replace_run_text(s[15], "Альтернатива:", "Альтернатива: ")  # restore trailing space
    replace_run_text(s[15], "Парето-frontier, ~70 мс", "~3× швидше за BGE-M3, ~70 мс")
    replace_run_text(s[15], "найвища nDCG, але повільний", "найвища nDCG, ~228 мс")
    # Replace outdated "Без GPU: Qwen3 / повільний, конкурентний"
    replace_run_text(s[15], "Без GPU:", "Лідер на Legal: ")
    replace_run_text(s[15], "повільний, конкурентний",
                     "nDCG@10=0.32, але ~2× повільніший")

    # ── Slide 18 — recommendations ─────────────────────────────────────
    print("Slide 18: recommendations — 13x -> 3x, GPU-only -> domain-specific")
    replace_run_text(
        s[17],
        "Прийнятна якість + ~13× швидше за BGE-M3",
        "Прийнятна якість + ~3× швидше за BGE-M3",
    )
    # First fix already collapsed "  —  лише з GPU" to "лідер на Legal" — restore separator
    replace_run_text(s[17], "лідер на Legal", "  —  лідер на юридичному домені")
    replace_run_text(
        s[17],
        "Виграє BGE-M3 за nDCG на юридичному домені; ~2× повільніший",
        "Найвища nDCG@10 на Legal (0.32), але ~2× повільніший за BGE-M3",
    )
    # Older path (in case fix is rerun on un-fixed pptx):
    replace_run_text(s[17], "Висока якість, але непрактична на CPU",
                     "Найвища nDCG@10 на Legal (0.32), але ~2× повільніший за BGE-M3")

    # ── Save ───────────────────────────────────────────────────────────
    prs.save(str(PPTX))
    print(f"\nSaved: {PPTX} ({PPTX.stat().st_size:,} bytes)")


if __name__ == "__main__":
    main()
