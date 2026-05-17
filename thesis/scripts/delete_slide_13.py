"""
Delete slide 13 (Результати: nDCG@10) — redundant with slide 14 detailed
metrics. The bootstrap CI block had wording the student already wanted
trimmed elsewhere («перевага семантичного пошуку над лексичним»).

Also updates page-number footers across remaining slides from «N / 18»
to «N / 17» (one renumbering for slides 1-12 unchanged, 14-18 shift
to 13-17).
"""
import sys
import re
from pathlib import Path
from pptx import Presentation

PPTX = Path(__file__).parent.parent / "2026_М_ПІ_ІПЗм-24-1_Нестеренко_В_В.pptx"


def delete_slide(prs, index_zero_based: int) -> None:
    """Remove slide at given 0-based index from the slide list."""
    xml_slides = prs.slides._sldIdLst
    slides = list(xml_slides)
    xml_slides.remove(slides[index_zero_based])


def update_page_footer(slide, old_total: int, new_total: int, new_index: int) -> bool:
    """Find any «NN / OLD_TOTAL» pattern and replace with «new_index / new_total».
    Returns True if a footer was found and updated.
    """
    pattern = re.compile(rf"(\d{{1,2}})\s*/\s*{old_total}")
    new_footer = f"{new_index:02d} / {new_total}"
    changed = False
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                if pattern.search(run.text):
                    run.text = pattern.sub(new_footer, run.text)
                    changed = True
                    return changed
            # Fallback: text spans multiple runs
            combined = "".join(r.text for r in para.runs)
            if pattern.search(combined):
                new_combined = pattern.sub(new_footer, combined)
                if para.runs:
                    para.runs[0].text = new_combined
                    for r in para.runs[1:]:
                        r.text = ""
                changed = True
                return changed
    return changed


def main() -> None:
    sys.stdout.reconfigure(encoding="utf-8")
    prs = Presentation(str(PPTX))

    before = len(prs.slides)
    print(f"Before: {before} slides")

    # Delete slide at 0-based index 12 (= slide 13 in 1-based)
    delete_slide(prs, 12)

    after = len(prs.slides)
    print(f"After:  {after} slides")
    assert after == before - 1, "Expected exactly one slide removed"

    # Update page-number footer on each remaining slide
    for new_idx, slide in enumerate(prs.slides, start=1):
        ok = update_page_footer(slide, old_total=18, new_total=after, new_index=new_idx)
        status = "OK  " if ok else "skip"
        print(f"  [{status}] slide now at position {new_idx}/{after}")

    prs.save(str(PPTX))
    print(f"\nSaved: {PPTX}")
    print("Next: re-run add_slide_time_intervals.py to refresh ⏱ markers")


if __name__ == "__main__":
    main()
