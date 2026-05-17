"""
Make BGE-M3 and BM25 subtitles on slide 17 (Підсумки) concrete with
numbers — supervisor flagged «якість, стабільність» as too abstract.
"""
import sys
from pathlib import Path
from pptx import Presentation

PPTX = Path(__file__).parent.parent / "2026_М_ПІ_ІПЗм-24-1_Нестеренко_В_В.pptx"

REPLACEMENTS = {
    "Найкраща якість, стабільність у різних доменах":
        "Лідер за nDCG@10 у tech (0.672) і medical (0.434), топ-2 у legal",
    "Залишається сильним у доменах, де важлива швидкість пошуку":
        "Найшвидша (1–3 мс/запит) — у >100× швидше за BGE-M3",
    # Re-running idempotently for the BM25 cell after first wrong number:
    "Найшвидша (~5 мс/запит) — у ~45× швидше за BGE-M3":
        "Найшвидша (1–3 мс/запит) — у >100× швидше за BGE-M3",
}


def main() -> None:
    sys.stdout.reconfigure(encoding="utf-8")
    prs = Presentation(str(PPTX))
    slide = prs.slides[16]  # slide 17 (Підсумки)
    changed = 0
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                if run.text in REPLACEMENTS:
                    print(f"  OLD: {run.text}")
                    run.text = REPLACEMENTS[run.text]
                    print(f"  NEW: {run.text}\n")
                    changed += 1
    # Fallback: paragraph-level replacement for split-run cases
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        for para in shape.text_frame.paragraphs:
            full = para.text
            if full in REPLACEMENTS:
                print(f"  OLD: {full}")
                new = REPLACEMENTS[full]
                if para.runs:
                    for r in para.runs[1:]:
                        r._r.getparent().remove(r._r)
                    para.runs[0].text = new
                print(f"  NEW: {new}\n")
                changed += 1

    print(f"Total replacements: {changed}")
    prs.save(str(PPTX))
    print(f"Saved: {PPTX}")


if __name__ == "__main__":
    main()
