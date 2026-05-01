"""
Заповнює Nesterenko_Presentation.pptx реальним змістом магістерського захисту.
Зберігає структуру шаблону (тема luxe, декоративні зображення, slide numbers).
Замінює текст у title + body placeholder'ах кожного зі 12 слайдів.
"""
import sys, copy, pathlib
from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor

sys.stdout.reconfigure(encoding='utf-8')

ROOT = pathlib.Path("D:/repos/kb-semantic-search-benchmark")
PPTX = ROOT / "Nesterenko_Presentation.pptx"

prs = Presentation(str(PPTX))


def find_title_and_body(slide):
    """Повертає (title_shape, body_shape). У шаблоні це 1-й і 2-й Google Shape."""
    google_shapes = [s for s in slide.shapes if s.has_text_frame and "Google" in s.name]
    if len(google_shapes) >= 2:
        return google_shapes[0], google_shapes[1]
    return None, None


def replace_text(shape, lines, title_size=None, body_size=None):
    """Замінює всі параграфи у shape новими рядками, зберігаючи стиль першого run."""
    if shape is None:
        return
    tf = shape.text_frame
    # Зберегти форматування першого run з першого параграфа
    first_para = tf.paragraphs[0]
    first_run = first_para.runs[0] if first_para.runs else None
    saved_size = first_run.font.size if first_run else None
    saved_name = first_run.font.name if first_run else None
    saved_bold = first_run.font.bold if first_run else None

    # Очистити всі параграфи
    p_elements = list(tf._txBody.findall(
        "{http://schemas.openxmlformats.org/drawingml/2006/main}p"
    ))
    for p in p_elements:
        tf._txBody.remove(p)

    # Додати нові параграфи
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0] if False else tf.add_paragraph()
        else:
            p = tf.add_paragraph()
        run = p.add_run()
        run.text = line
        # Відновлюємо форматування
        if saved_size:
            run.font.size = saved_size
        if saved_name:
            run.font.name = saved_name
        if saved_bold is not None:
            run.font.bold = saved_bold


# ── Slide 1 (TITLE) — лише виправляємо "звання, посада" ────────────────
slide1 = prs.slides[0]
for shape in slide1.shapes:
    if shape.has_text_frame and "Google Shape;63" in shape.name:
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                if "звання, посада" in (run.text or ""):
                    run.text = run.text.replace(
                        "звання, посада ",
                        "к.т.н., доцент "
                    )

# ── Slide 2: Дослідження → Актуальність + мета + задачі ────────────────
s2_title, s2_body = find_title_and_body(prs.slides[1])
replace_text(s2_title, ["Актуальність та мета дослідження"])
replace_text(s2_body, [
    "Актуальність: експоненційне зростання обсягів цифрової інформації у корпоративних базах знань; класичний keyword-пошук обмежений; розвиток LLM/RAG-систем робить retrieval-якість критичною",
    "Об'єкт: процеси формування та використання векторних подань текстових даних у системах семантичного пошуку",
    "Предмет: моделі векторних ембеддінгів для корпоративних баз знань",
    "Мета: визначення доцільності застосування сучасних моделей векторних ембеддінгів для підвищення ефективності семантичного пошуку у корпоративних базах знань шляхом їх порівняльного експериментального дослідження",
    "Задачі: огляд предметної галузі; аналіз сучасних embedding-моделей; формування benchmark dataset; реалізація системи; експериментальна оцінка; багатокритеріальний вибір оптимальної моделі",
])

# ── Slide 3: Огляд літератури (аналогів) ───────────────────────────────
s3_title, s3_body = find_title_and_body(prs.slides[2])
replace_text(s3_title, ["Огляд літератури"])
replace_text(s3_body, [
    "Фундаментальні роботи: Mikolov (Word2Vec), Pennington (GloVe), Devlin (BERT), Reimers (Sentence-BERT)",
    "Сучасні embedding-моделі: Wang (E5), Chen (BGE-M3), Nussbaum (nomic-embed-text), Zhang (Qwen3-Embedding)",
    "Лексичний базелайн: Robertson, Zaragoza (BM25)",
    "Бенчмарки і метрики: Thakur (BEIR), Muennighoff (MTEB), Järvelin (nDCG)",
    "RAG-парадигма: Lewis (Retrieval-Augmented Generation), Gao (RAG Survey)",
    "Прогалина: відсутність порівняльних досліджень embedding-моделей на україномовних доменно-специфічних колекціях",
    "Загалом використано 32 наукових джерела",
])

# ── Slide 4: Постановка задачі ─────────────────────────────────────────
s4_title, s4_body = find_title_and_body(prs.slides[3])
replace_text(s4_title, ["Постановка задачі"])
replace_text(s4_body, [
    "Проблема: у реальних корпоративних системах недостатньо просто застосувати сучасну embedding-модель — необхідно обґрунтовано визначити оптимальну за балансом якості, швидкодії та ресурсів",
    "Очікувані результати:",
    "  – реалізована система семантичного пошуку з підтримкою документів PDF/DOCX/TXT/MD",
    "  – власноруч сформований domain-specific україномовний benchmark dataset",
    "  – кількісне порівняння 4 embedding-моделей + BM25 baseline на 3 доменах",
    "  – багатокритеріальний вибір оптимальної моделі (Парето + лінійна згортка)",
    "  – обґрунтовані практичні рекомендації для впровадження",
])

# ── Slide 5: Методологія ───────────────────────────────────────────────
s5_title, s5_body = find_title_and_body(prs.slides[4])
replace_text(s5_title, ["Методологія дослідження"])
replace_text(s5_body, [
    "Метрики оцінювання retrieval-якості: nDCG@10, MRR@10, Recall@10, P@10",
    "Метрика швидкодії: середній час відповіді на запит (мс/запит)",
    "Статистичний аналіз: bootstrap-оцінка 95% довірчого інтервалу для nDCG@10 (n = 2000)",
    "Багатокритеріальний вибір: принцип Парето-домінування + лінійна адитивна згортка з нормалізацією критеріїв",
    "Метод призначення ваг: ранжування критеріїв за важливістю з доменно-специфічними профілями",
    "Інструментарій: Python 3.11, sentence-transformers, FAISS, rank_bm25, Flask, NumPy, scikit-learn",
])

# ── Slide 6: Архітектура системи ───────────────────────────────────────
s6_title, s6_body = find_title_and_body(prs.slides[5])
replace_text(s6_title, ["Архітектура системи"])
replace_text(s6_body, [
    "Pipeline експерименту: документи → попередня обробка → чанкінг → ембеддінг → FAISS-індекс → запит → top-k результати → метрики",
    "Основні компоненти:",
    "  – build_index.py — побудова індексу для (модель, домен)",
    "  – evaluate_benchmark.py — обчислення метрик retrieval-якості",
    "  – embedding_models.py — уніфікований інтерфейс для 5 моделей",
    "  – run_all_benchmarks.py — масовий прогін усіх (модель, домен)",
    "Web-інтерфейс: Flask-додаток з вибором моделі, домену та інтерактивним пошуком",
    "Зберігання артефактів: faiss.index + meta.jsonl + model.json для кожної (модель, домен)",
])

# ── Slide 7: ПЗ ────────────────────────────────────────────────────────
s7_title, s7_body = find_title_and_body(prs.slides[6])
replace_text(s7_title, ["Програмне забезпечення"])
replace_text(s7_body, [
    "Мова та середовище: Python 3.11, віртуальне середовище venv",
    "Embedding-фреймворки: sentence-transformers (Hugging Face), trust_remote_code для nomic та Qwen3",
    "Векторний індекс: FAISS IndexFlatIP з L2-нормалізацією векторів",
    "Лексичний базелайн: rank_bm25 (Okapi BM25)",
    "Web-фреймворк: Flask (Python)",
    "Підтримувані формати документів: PDF, DOCX, TXT, MD",
    "Моделі: BGE-M3 (BAAI), E5-base (intfloat/multilingual-e5-base), nomic-embed-text-v1.5, Qwen3-Embedding-0.6B (Alibaba)",
])

# ── Slide 8: Зміст експерименту (KEY DATASET SLIDE) ────────────────────
s8_title, s8_body = find_title_and_body(prs.slides[7])
replace_text(s8_title, ["Benchmark dataset та умови експерименту"])
replace_text(s8_body, [
    "Власноруч сформований domain-specific україномовний benchmark dataset",
    "Технічний домен: 10 документів, 178 чанків (програмна інженерія, ML, RAG, аномалії, 3D-друк)",
    "Юридичний домен: 13 документів, 1046 чанків (кодекси, нормативно-правові акти)",
    "Медичний домен: 9 документів, 267 чанків (клінічні протоколи, медичні рекомендації)",
    "Усього: 32 документи, 1 491 чанк, 300 пошукових запитів (по 100 на домен)",
    "Релевантні відповідності (qrels): сформовані автором вручну на рівні чанків",
    "Учасники експерименту: 4 нейронні embedding-моделі + BM25 базелайн на кожному з 3 доменів (15 експериментів)",
])

# ── Slide 9: Результати експерименту ───────────────────────────────────
s9_title, s9_body = find_title_and_body(prs.slides[8])
replace_text(s9_title, ["Результати експерименту: nDCG@10"])
replace_text(s9_body, [
    "Зведена таблиця nDCG@10 (вище — краще):",
    "BGE-M3: технічний 0.6722  |  юридичний 0.3065  |  медичний 0.4339",
    "Qwen3-Embedding-0.6B: 0.6325  |  0.3199  |  0.3629",
    "E5-base: 0.6121  |  0.2567  |  0.3909",
    "BM25 (базелайн): 0.4861  |  0.1875  |  0.3222",
    "nomic-embed-text-v1.5: 0.3765  |  0.0951  |  0.1668",
    "Bootstrap 95% CI (n = 2000) для технічного домену: BGE-M3 [0.606, 0.737] vs BM25 [0.431, 0.541] — інтервали не перетинаються (статистично значуща перевага)",
])

# ── Slide 10: Аналіз результатів ───────────────────────────────────────
s10_title, s10_body = find_title_and_body(prs.slides[9])
replace_text(s10_title, ["Аналіз результатів"])
replace_text(s10_body, [
    "BGE-M3 — лідер у всіх 3 доменах за nDCG@10, MRR@10 та Recall@10",
    "Парето-оптимальні альтернативи: P = {BGE-M3, E5-base}",
    "Інтегральна оцінка U (доменно-специфічні ваги): BGE-M3 → 0.920 (тех.) / 0.927 (юр.) / 0.901 (мед.)",
    "E5-base — швидка альтернатива (~225 мс/запит проти ~2960 мс у BGE-M3)",
    "Qwen3-Embedding-0.6B — конкурентна якість, але ~1500 мс/запит на CPU робить її непрактичною без GPU",
    "nomic-embed-text-v1.5 — найслабші результати; в окремих доменах поступається BM25",
    "Метою досягнуто: визначено доцільність застосування сучасних embedding-моделей; обрано BGE-M3 як рекомендовану модель",
])

# ── Slide 11: Публікація ───────────────────────────────────────────────
s11_title, s11_body = find_title_and_body(prs.slides[10])
replace_text(s11_title, ["Апробація результатів"])
replace_text(s11_body, [
    "Результати дослідження апробовано на двох наукових конференціях:",
    "1. 2-я Міжнародна науково-практична конференція «Innovative Research in Science and Economy» — стаття «Порівняльний аналіз ефективності моделей векторних ембеддінгів для задач семантичного пошуку в корпоративних базах знань»",
    "2. 30-й Міжнародний молодіжний форум «Радіоелектроніка та молодь у XXI столітті» (ХНУРЕ) — тези «Підвищення якості семантичного пошуку в базах знань із використанням векторних подань»",
    "Опубліковано 2 наукові праці у співавторстві з науковим керівником",
])

# ── Slide 12: Підсумки ─────────────────────────────────────────────────
s12_title, s12_body = find_title_and_body(prs.slides[11])
replace_text(s12_title, ["Підсумки та перспективи"])
replace_text(s12_body, [
    "Поставлену мету досягнуто: визначено доцільність застосування сучасних моделей векторних ембеддінгів для семантичного пошуку у корпоративних базах знань",
    "Розроблено програмну систему та власноруч сформовано benchmark dataset для україномовних доменно-специфічних колекцій",
    "Практичні рекомендації для впровадження:",
    "  – BGE-M3 як основна модель за якістю та стабільністю",
    "  – E5-base як швидка альтернатива для інтерактивного пошуку",
    "  – Qwen3-Embedding — лише за наявності GPU",
    "  – BM25 залишається конкурентним базелайном у певних доменах",
    "Перспективи: розширення benchmark-наборів, інтеграція reranking-підходів, повноцінна RAG-архітектура",
    "Дякую за увагу!",
])

prs.save(str(PPTX))
print(f"Saved: {PPTX} ({PPTX.stat().st_size:,} bytes)")
print(f"Slides: {len(prs.slides)}")
