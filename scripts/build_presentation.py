"""
Комплексна побудова презентації магістерського захисту.
17 слайдів, фокус на дослідженнях.

Структура:
  1.  Title
  2.  Актуальність + мета + задачі
  3.  Огляд літератури — timeline
  4.  Постановка задачі
  5.  [НОВЕ] Аналіз 4 embedding-моделей
  6.  Методологія дослідження
  7.  [НОВЕ] Метрики оцінювання — формули
  8.  [НОВЕ] Багатокритеріальний вибір — методологія
  9.  Архітектура системи
  10. Програмне забезпечення
  11. Benchmark dataset
  12. Результати: nDCG@10
  13. [НОВЕ] Деталізовані метрики (MRR, Recall, P, ms)
  14. МКВ: Парето + інтегральні оцінки
  15. [НОВЕ] Якість vs швидкодія (scatter)
  16. Апробація
  17. Підсумки + рекомендації
"""
import sys, copy, pathlib
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR, MSO_AUTO_SIZE
from pptx.dml.color import RGBColor
from pptx.chart.data import CategoryChartData, XyChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION, XL_LABEL_POSITION
from pptx.oxml.ns import qn
from lxml import etree

sys.stdout.reconfigure(encoding='utf-8')

ROOT = pathlib.Path("D:/repos/kb-semantic-search-benchmark")
PPTX = ROOT / "Nesterenko_Presentation.pptx"

# ── Кольорова палітра ─────────────────────────────────────────────────
PRIMARY    = RGBColor(0x1E, 0x3A, 0x8A)
PRIMARY_LT = RGBColor(0x3B, 0x82, 0xF6)
ACCENT     = RGBColor(0xF5, 0x9E, 0x0B)
SUCCESS    = RGBColor(0x16, 0xA3, 0x4A)
DANGER     = RGBColor(0xDC, 0x26, 0x26)
GRAY_DARK  = RGBColor(0x1F, 0x29, 0x37)
GRAY_MED   = RGBColor(0x64, 0x74, 0x8B)
GRAY_LIGHT = RGBColor(0xF1, 0xF5, 0xF9)
GRAY_BORD  = RGBColor(0xCB, 0xD5, 0xE1)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
TEAL       = RGBColor(0x0D, 0x94, 0x88)
PURPLE     = RGBColor(0x7C, 0x3A, 0xED)

DOMAIN_TECH  = RGBColor(0x1E, 0x40, 0xAF)
DOMAIN_LEGAL = RGBColor(0xB4, 0x53, 0x09)
DOMAIN_MED   = RGBColor(0xBE, 0x18, 0x5D)


def IN(v): return Inches(v)


def remove_shape(shape):
    sp = shape._element
    sp.getparent().remove(sp)


def find_shape_by_name(slide, substr):
    for s in slide.shapes:
        if substr in s.name:
            return s
    return None


def remove_body_placeholder(slide):
    google_shapes = [s for s in slide.shapes if "Google" in s.name and s.has_text_frame]
    if len(google_shapes) >= 2:
        remove_shape(google_shapes[1])


def set_title(slide, new_title, font_size=24):
    google_shapes = [s for s in slide.shapes if "Google" in s.name and s.has_text_frame]
    if google_shapes:
        title = google_shapes[0]
        tf = title.text_frame
        try:
            tf.auto_size = MSO_AUTO_SIZE.NONE
        except Exception:
            pass
        p_elements = list(tf._txBody.findall(qn("a:p")))
        for p in p_elements[1:]:
            tf._txBody.remove(p)
        first_p = tf.paragraphs[0]
        for r in list(first_p._p.findall(qn("a:r"))):
            first_p._p.remove(r)
        run = first_p.add_run()
        run.text = new_title
        run.font.size = Pt(font_size)
        run.font.bold = True
        run.font.color.rgb = GRAY_DARK
        run.font.name = "Calibri"


def set_slide_number(slide, num):
    """Updates the slide number TextBox."""
    for shape in slide.shapes:
        if shape.name == "TextBox 1" and shape.has_text_frame:
            tf = shape.text_frame
            for p in tf.paragraphs:
                for r in p.runs:
                    r.text = str(num)
                    return


def add_text(slide, text, left, top, width, height,
             font_size=14, font_name="Calibri", bold=False, italic=False,
             color=None, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
             fill=None, line_color=None, line_width=None):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = IN(0.05)
    tf.margin_right = IN(0.05)
    tf.margin_top = IN(0.05)
    tf.margin_bottom = IN(0.05)
    tf.vertical_anchor = anchor
    if fill:
        tb.fill.solid()
        tb.fill.fore_color.rgb = fill
    if line_color:
        tb.line.color.rgb = line_color
        if line_width:
            tb.line.width = line_width
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.name = font_name
    run.font.bold = bold
    run.font.italic = italic
    if color:
        run.font.color.rgb = color
    return tb


def add_shape(slide, shape_type, left, top, width, height,
              fill=None, line_color=None, line_width=None,
              text=None, font_size=14, bold=False, color=None,
              align=PP_ALIGN.CENTER, font_name="Calibri",
              vertical_anchor=MSO_ANCHOR.MIDDLE):
    sh = slide.shapes.add_shape(shape_type, left, top, width, height)
    if fill is not None:
        if fill == "none":
            sh.fill.background()
        else:
            sh.fill.solid()
            sh.fill.fore_color.rgb = fill
    if line_color:
        sh.line.color.rgb = line_color
        if line_width:
            sh.line.width = line_width
    elif line_color is None:
        sh.line.fill.background()
    if text is not None:
        tf = sh.text_frame
        tf.word_wrap = True
        tf.margin_left = IN(0.06)
        tf.margin_right = IN(0.06)
        tf.margin_top = IN(0.04)
        tf.margin_bottom = IN(0.04)
        tf.vertical_anchor = vertical_anchor
        p = tf.paragraphs[0]
        p.alignment = align
        run = p.add_run()
        run.text = text
        run.font.size = Pt(font_size)
        run.font.name = font_name
        run.font.bold = bold
        if color:
            run.font.color.rgb = color
    return sh


def add_card(slide, left, top, width, height, title, body_lines,
             title_bg=PRIMARY, title_color=WHITE, body_bg=GRAY_LIGHT,
             title_size=12, body_size=10, body_color=GRAY_DARK):
    title_h = IN(0.35)
    add_shape(slide, MSO_SHAPE.RECTANGLE, left, top, width, title_h,
              fill=title_bg, line_color=None,
              text=title, font_size=title_size, bold=True,
              color=title_color, align=PP_ALIGN.CENTER)
    body = add_shape(slide, MSO_SHAPE.RECTANGLE, left, top + title_h, width, height - title_h,
                     fill=body_bg, line_color=GRAY_BORD)
    tf = body.text_frame
    tf.word_wrap = True
    tf.margin_left = IN(0.1); tf.margin_right = IN(0.1)
    tf.margin_top = IN(0.08); tf.margin_bottom = IN(0.08)
    p = tf.paragraphs[0]
    for i, line in enumerate(body_lines):
        if i == 0:
            para = p
        else:
            para = tf.add_paragraph()
        para.alignment = PP_ALIGN.LEFT
        run = para.add_run()
        run.text = line
        run.font.size = Pt(body_size)
        run.font.name = "Calibri"
        run.font.color.rgb = body_color


def add_arrow(slide, x1, y1, x2, y2, color=PRIMARY, weight=2.0):
    line = slide.shapes.add_connector(1, x1, y1, x2, y2)
    line.line.color.rgb = color
    line.line.width = Pt(weight)
    el = line.line._get_or_add_ln()
    tail = etree.SubElement(el, qn("a:tailEnd"))
    tail.set("type", "triangle")
    tail.set("w", "med")
    tail.set("h", "med")
    return line


def duplicate_slide(prs, src_slide_idx):
    """Дублює слайд (з копією зображень/relationships) в кінець."""
    src = prs.slides[src_slide_idx]
    new_slide = prs.slides.add_slide(src.slide_layout)

    # Видалити всі shape, успадковані від layout
    for shape in list(new_slide.shapes):
        sp = shape.element
        sp.getparent().remove(sp)

    # Скопіювати relationships з джерела (для зображень)
    rel_id_map = {}  # old_rId → new_rId
    for rel_id, rel in src.part.rels.items():
        if rel.is_external:
            continue
        if 'image' in rel.reltype or 'picture' in rel.reltype:
            new_rel_id = new_slide.part.relate_to(rel.target_part, rel.reltype)
            rel_id_map[rel_id] = new_rel_id

    # Скопіювати XML усіх shape, оновлюючи rId де потрібно
    for shape in src.shapes:
        new_el = copy.deepcopy(shape.element)
        # Замінити rId у всіх атрибутах
        for el in new_el.iter():
            for attr_name in (qn("r:embed"), qn("r:link")):
                old_rid = el.get(attr_name)
                if old_rid and old_rid in rel_id_map:
                    el.set(attr_name, rel_id_map[old_rid])
        new_slide.shapes._spTree.append(new_el)

    return new_slide


def reorder_slides(prs, new_order):
    """Реорганізує слайди згідно з new_order (список 0-based індексів)."""
    sldIdLst = prs.slides._sldIdLst
    sld_ids = list(sldIdLst)
    # Замість видалення — переміщуємо XML елементи
    for sld_id in sld_ids:
        sldIdLst.remove(sld_id)
    for new_idx in new_order:
        sldIdLst.append(sld_ids[new_idx])


def clear_slide_content(slide, keep_picture=True, keep_textbox=True):
    """Видаляє всі shape окрім зображень і TextBox 1 (slide number)."""
    for shape in list(slide.shapes):
        if "Google" in shape.name and "PICTURE" in str(shape.shape_type):
            if keep_picture:
                continue
        if shape.name == "TextBox 1":
            if keep_textbox:
                continue
        # Зберігаємо титульний placeholder (перший Google Shape з text frame)
        # Видаляємо все інше
        try:
            sp = shape.element
            sp.getparent().remove(sp)
        except:
            pass


# ════════════════════════════════════════════════════════════════════════
# Зона контенту: x=0.34", y=1.34", w=9.32", h=3.55"
CX = IN(0.34)
CY = IN(1.34)
CW = IN(9.32)
CH = IN(3.55)


# ── ВМІСТ СЛАЙДІВ ─────────────────────────────────────────────────────

def fill_slide_title(slide):
    """Slide 1: Title — лише фікс 'звання, посада'."""
    for shape in slide.shapes:
        if shape.has_text_frame and "Google Shape;63" in shape.name:
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    if "звання, посада" in (run.text or ""):
                        run.text = run.text.replace("звання, посада ", "к.т.н., доцент ")


def fill_slide_actuality(slide):
    """Slide 2: Актуальність, мета, задачі."""
    set_title(slide, "Актуальність, мета та задачі")
    remove_body_placeholder(slide)

    # Goal hero block
    goal_h = IN(1.0)
    goal_box = add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, CX, CY, CW, goal_h,
                         fill=PRIMARY, line_color=None)
    tf = goal_box.text_frame
    tf.word_wrap = True
    tf.margin_left = IN(0.18); tf.margin_right = IN(0.18)
    tf.margin_top = IN(0.08); tf.margin_bottom = IN(0.08)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    r = p.add_run(); r.text = "МЕТА: "
    r.font.size = Pt(13); r.font.bold = True; r.font.color.rgb = ACCENT; r.font.name = "Calibri"
    r = p.add_run()
    r.text = ("визначення доцільності застосування сучасних моделей векторних "
              "ембеддінгів для підвищення ефективності семантичного пошуку текстових "
              "документів у корпоративних базах знань шляхом їх порівняльного "
              "експериментального дослідження")
    r.font.size = Pt(11); r.font.color.rgb = WHITE; r.font.name = "Calibri"

    # 6 task cards
    tasks = [
        ("1", "Аналіз\nпредметної галузі"),
        ("2", "Огляд\nembedding-моделей"),
        ("3", "Формування\nbenchmark dataset"),
        ("4", "Реалізація\nсистеми пошуку"),
        ("5", "Оцінка\nякості моделей"),
        ("6", "Вибір\nоптимальної моделі"),
    ]
    tasks_top = CY + goal_h + IN(0.18)
    card_w = IN(1.45); card_h = IN(0.95); gap = IN(0.08)
    total_w = card_w * 6 + gap * 5
    start_x = CX + (CW - total_w) // 2
    for i, (num, text) in enumerate(tasks):
        x = start_x + i * (card_w + gap)
        # Card first
        box = add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, x, tasks_top, card_w, card_h,
                        fill=GRAY_LIGHT, line_color=GRAY_BORD, line_width=Pt(0.75))
        tf = box.text_frame
        tf.word_wrap = True
        tf.margin_left = IN(0.4); tf.margin_right = IN(0.06)
        tf.margin_top = IN(0.06); tf.margin_bottom = IN(0.06)
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]; p.alignment = PP_ALIGN.LEFT
        r = p.add_run(); r.text = text
        r.font.size = Pt(9.5); r.font.color.rgb = GRAY_DARK; r.font.name = "Calibri"; r.font.bold = True
        # Number circle
        circle_d = IN(0.32)
        add_shape(slide, MSO_SHAPE.OVAL, x + IN(0.06), tasks_top + IN(0.05), circle_d, circle_d,
                  fill=ACCENT, line_color=None,
                  text=num, font_size=12, bold=True, color=WHITE)

    # Об'єкт + Предмет
    op_top = tasks_top + card_h + IN(0.18)
    op_h = IN(0.85)
    op_w = (CW - IN(0.15)) // 2
    add_card(slide, CX, op_top, op_w, op_h,
             title="ОБ'ЄКТ ДОСЛІДЖЕННЯ", title_bg=GRAY_DARK, title_size=9,
             body_lines=["процеси формування та використання векторних подань текстових даних",
                         "у системах семантичного пошуку"],
             body_size=9.5)
    add_card(slide, CX + op_w + IN(0.15), op_top, op_w, op_h,
             title="ПРЕДМЕТ ДОСЛІДЖЕННЯ", title_bg=GRAY_DARK, title_size=9,
             body_lines=["моделі векторних ембеддінгів для семантичного пошуку",
                         "у корпоративних базах знань"],
             body_size=9.5)


def fill_slide_timeline(slide):
    """Slide 3: Огляд літератури — timeline."""
    set_title(slide, "Огляд літератури: еволюція моделей")
    remove_body_placeholder(slide)

    tl_y = CY + IN(1.4)
    tl_x1 = CX + IN(0.4)
    tl_x2 = CX + CW - IN(0.4)
    line = slide.shapes.add_connector(1, tl_x1, tl_y, tl_x2, tl_y)
    line.line.color.rgb = PRIMARY; line.line.width = Pt(3)

    milestones = [
        ("2013", "Word2Vec", "Mikolov", "Статистичні\nвекторні подання"),
        ("2018", "BERT", "Devlin", "Трансформери,\nконтекстність"),
        ("2019", "Sentence-BERT", "Reimers", "Семантичні\nвекторизатори"),
        ("2022", "E5", "Wang", "Multilingual,\ncontrastive"),
        ("2024", "BGE-M3 / nomic", "Chen / Nussbaum", "Multi-functionality,\nlong context"),
        ("2025", "Qwen3-Embedding", "Zhang", "Decoder-based,\ninstruction-aware"),
    ]
    ms_count = len(milestones)
    seg_w = (tl_x2 - tl_x1) / (ms_count - 1)
    for i, (year, name, author, desc) in enumerate(milestones):
        cx = tl_x1 + int(seg_w * i)
        dot_d = IN(0.22)
        add_shape(slide, MSO_SHAPE.OVAL, cx - dot_d // 2, tl_y - dot_d // 2, dot_d, dot_d,
                  fill=ACCENT, line_color=WHITE, line_width=Pt(2))
        add_text(slide, year, cx - IN(0.4), tl_y - IN(0.85), IN(0.8), IN(0.25),
                 font_size=11, bold=True, color=PRIMARY, align=PP_ALIGN.CENTER)
        add_text(slide, name, cx - IN(0.65), tl_y - IN(0.58), IN(1.3), IN(0.22),
                 font_size=10, bold=True, color=GRAY_DARK, align=PP_ALIGN.CENTER)
        add_text(slide, author, cx - IN(0.65), tl_y + IN(0.1), IN(1.3), IN(0.22),
                 font_size=8, italic=True, color=GRAY_MED, align=PP_ALIGN.CENTER)
        add_text(slide, desc, cx - IN(0.65), tl_y + IN(0.32), IN(1.3), IN(0.5),
                 font_size=8, color=GRAY_DARK, align=PP_ALIGN.CENTER)

    gap_top = CY + IN(2.85)
    gap_box = add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, CX, gap_top, CW, IN(0.65),
                        fill=GRAY_LIGHT, line_color=ACCENT, line_width=Pt(1.5))
    tf = gap_box.text_frame; tf.word_wrap = True
    tf.margin_left = IN(0.15); tf.margin_right = IN(0.15)
    tf.margin_top = IN(0.06); tf.margin_bottom = IN(0.06)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.LEFT
    r = p.add_run(); r.text = "ПРОГАЛИНА: "
    r.font.size = Pt(11); r.font.bold = True; r.font.color.rgb = ACCENT; r.font.name = "Calibri"
    r = p.add_run()
    r.text = ("відсутність порівняльних досліджень embedding-моделей на україномовних "
              "доменно-специфічних колекціях.   ▸ використано 32 наукових джерела "
              "(Lewis 2020 — RAG, Robertson 2009 — BM25, Thakur 2021 — BEIR, Muennighoff 2023 — MTEB)")
    r.font.size = Pt(10); r.font.color.rgb = GRAY_DARK; r.font.name = "Calibri"


def fill_slide_problem(slide):
    """Slide 4: Постановка задачі."""
    set_title(slide, "Постановка задачі")
    remove_body_placeholder(slide)

    lw = IN(4.5); rw = CW - lw - IN(0.2)
    prob_box = add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, CX, CY, lw, CH,
                         fill=PRIMARY, line_color=None)
    tf = prob_box.text_frame; tf.word_wrap = True
    tf.margin_left = IN(0.2); tf.margin_right = IN(0.2)
    tf.margin_top = IN(0.15); tf.margin_bottom = IN(0.15)
    tf.vertical_anchor = MSO_ANCHOR.TOP
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.LEFT
    r = p.add_run(); r.text = "ПРОБЛЕМА"
    r.font.size = Pt(15); r.font.bold = True; r.font.color.rgb = ACCENT; r.font.name = "Calibri"
    p = tf.add_paragraph(); p.alignment = PP_ALIGN.LEFT; p.space_before = Pt(8)
    r = p.add_run()
    r.text = ("У реальних корпоративних системах недостатньо просто застосувати сучасну "
              "embedding-модель. Необхідно обґрунтовано визначити модель, яка забезпечує "
              "оптимальний баланс між:")
    r.font.size = Pt(11); r.font.color.rgb = WHITE; r.font.name = "Calibri"
    for crit in ["• якістю семантичного пошуку", "• швидкодією на CPU/GPU",
                 "• вимогами до пам'яті та ресурсів", "• мовною підтримкою", "• типом колекції документів"]:
        p = tf.add_paragraph(); p.alignment = PP_ALIGN.LEFT
        r = p.add_run(); r.text = crit
        r.font.size = Pt(10.5); r.font.color.rgb = WHITE; r.font.name = "Calibri"

    right_x = CX + lw + IN(0.2)
    add_text(slide, "ОЧІКУВАНІ РЕЗУЛЬТАТИ", right_x, CY, rw, IN(0.3),
             font_size=12, bold=True, color=PRIMARY, align=PP_ALIGN.LEFT)
    outcomes = [
        "Реалізована система семантичного пошуку (PDF/DOCX/TXT/MD)",
        "Власноруч сформований UA benchmark dataset (3 домени)",
        "Кількісне порівняння 4 embedding-моделей + BM25 baseline",
        "Багатокритеріальний вибір (Парето + лінійна згортка)",
        "Обґрунтовані практичні рекомендації для впровадження",
    ]
    out_top = CY + IN(0.35); out_h = IN(0.55); out_gap = IN(0.06)
    for i, txt in enumerate(outcomes):
        y = out_top + i * (out_h + out_gap)
        add_shape(slide, MSO_SHAPE.OVAL, right_x, y + IN(0.07), IN(0.4), IN(0.4),
                  fill=ACCENT, line_color=None,
                  text=str(i + 1), font_size=14, bold=True, color=WHITE)
        add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, right_x + IN(0.5), y, rw - IN(0.5), out_h,
                  fill=GRAY_LIGHT, line_color=GRAY_BORD,
                  text=txt, font_size=10, bold=False, color=GRAY_DARK,
                  align=PP_ALIGN.LEFT)


def fill_slide_models(slide):
    """Slide 5 [НОВЕ]: Аналіз 4 embedding-моделей."""
    set_title(slide, "Сучасні моделі векторних ембеддінгів")
    remove_body_placeholder(slide)

    # 2x2 grid of model cards
    models = [
        ("BGE-M3 (BAAI)", DOMAIN_TECH,
         [("Архітектура", "Encoder (трансформер)"),
          ("Параметри", "568 M"),
          ("Max seq length", "8 192 токенів"),
          ("Особливості", "Multi-functionality:\ndense + sparse + multi-vector"),
          ("Метод навчання", "Self-knowledge distillation")]),
        ("E5-base (intfloat)", PRIMARY_LT,
         [("Архітектура", "Encoder (multilingual-e5-base)"),
          ("Параметри", "278 M"),
          ("Max seq length", "512 токенів"),
          ("Особливості", "Префікси query: / passage: \nдля contrastive pretraining"),
          ("Метод навчання", "Weakly-supervised contrastive\n(1B+ текстових пар)")]),
        ("nomic-embed-text-v1.5", TEAL,
         [("Архітектура", "Encoder"),
          ("Параметри", "137 M"),
          ("Max seq length", "8 192 токенів"),
          ("Особливості", "Matryoshka representations\n(адаптивна розмірність)"),
          ("Метод навчання", "Long-context contrastive,\nповна reproducibility")]),
        ("Qwen3-Embedding-0.6B", PURPLE,
         [("Архітектура", "Decoder (LLM-based)"),
          ("Параметри", "596 M"),
          ("Max seq length", "32 768 токенів"),
          ("Особливості", "Instruction-aware,\nкерування через prompt"),
          ("Метод навчання", "Foundation model\nfine-tuning")]),
    ]
    card_w = (CW - IN(0.2)) // 2
    card_h = (CH - IN(0.15)) // 2
    for i, (name, color, fields) in enumerate(models):
        row = i // 2; col = i % 2
        x = CX + col * (card_w + IN(0.2))
        y = CY + row * (card_h + IN(0.15))
        # Header bar
        add_shape(slide, MSO_SHAPE.RECTANGLE, x, y, card_w, IN(0.4),
                  fill=color, line_color=None,
                  text=name, font_size=13, bold=True, color=WHITE)
        # Body
        body = add_shape(slide, MSO_SHAPE.RECTANGLE, x, y + IN(0.4), card_w, card_h - IN(0.4),
                         fill=GRAY_LIGHT, line_color=GRAY_BORD)
        tf = body.text_frame; tf.word_wrap = True
        tf.margin_left = IN(0.12); tf.margin_right = IN(0.1)
        tf.margin_top = IN(0.08); tf.margin_bottom = IN(0.05)
        tf.vertical_anchor = MSO_ANCHOR.TOP
        p = tf.paragraphs[0]
        for j, (key, val) in enumerate(fields):
            if j == 0:
                para = p
            else:
                para = tf.add_paragraph()
            para.alignment = PP_ALIGN.LEFT
            para.space_after = Pt(2)
            r = para.add_run(); r.text = key + ": "
            r.font.size = Pt(8.5); r.font.bold = True; r.font.color.rgb = color; r.font.name = "Calibri"
            r = para.add_run(); r.text = val
            r.font.size = Pt(8.5); r.font.color.rgb = GRAY_DARK; r.font.name = "Calibri"


def fill_slide_methodology(slide):
    """Slide 6: Методологія дослідження."""
    set_title(slide, "Методологія дослідження")
    remove_body_placeholder(slide)

    add_text(slide, "МЕТРИКИ ОЦІНЮВАННЯ RETRIEVAL-ЯКОСТІ", CX, CY, CW, IN(0.3),
             font_size=11, bold=True, color=PRIMARY, align=PP_ALIGN.LEFT)
    metrics = [
        ("nDCG@10", "Якість\nранжування", "(основна)", PRIMARY),
        ("MRR@10", "Позиція першого\nрелевантного", "", PRIMARY_LT),
        ("Recall@10", "Повнота\nпошуку", "", TEAL),
        ("P@10", "Точність\nвидачі", "", PURPLE),
        ("ms/query", "Час відповіді\n(швидкодія)", "", ACCENT),
    ]
    m_top = CY + IN(0.3); m_h = IN(1.05)
    m_w = (CW - IN(0.3)) // 5
    for i, (name, desc, note, color) in enumerate(metrics):
        x = CX + i * (m_w + IN(0.075))
        add_shape(slide, MSO_SHAPE.RECTANGLE, x, m_top, m_w, IN(0.42),
                  fill=color, line_color=None,
                  text=name, font_size=14, bold=True, color=WHITE)
        body = add_shape(slide, MSO_SHAPE.RECTANGLE, x, m_top + IN(0.42), m_w, m_h - IN(0.42),
                         fill=GRAY_LIGHT, line_color=GRAY_BORD)
        tf = body.text_frame; tf.word_wrap = True
        tf.margin_left = IN(0.08); tf.margin_right = IN(0.08)
        tf.margin_top = IN(0.06); tf.margin_bottom = IN(0.06)
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        r = p.add_run(); r.text = desc
        r.font.size = Pt(9.5); r.font.bold = True; r.font.color.rgb = GRAY_DARK; r.font.name = "Calibri"
        if note:
            p = tf.add_paragraph(); p.alignment = PP_ALIGN.CENTER
            r = p.add_run(); r.text = note
            r.font.size = Pt(8); r.font.italic = True; r.font.color.rgb = ACCENT; r.font.name = "Calibri"

    m2_top = m_top + m_h + IN(0.2)
    add_text(slide, "МЕТОДИ АНАЛІЗУ", CX, m2_top - IN(0.27), CW, IN(0.25),
             font_size=11, bold=True, color=PRIMARY, align=PP_ALIGN.LEFT)
    methods = [
        ("Bootstrap CI", "95% довірчий інтервал для nDCG@10\n(n = 2 000 повторень)"),
        ("Парето-домінування", "Виявлення Парето-оптимальних\nальтернатив у багатокритеріальному просторі"),
        ("Лінійна адитивна згортка", "Інтегральна оцінка U(aᵢ) з ваговими\nкоефіцієнтами критеріїв"),
    ]
    m2_h = IN(1.0); m2_w = (CW - IN(0.3)) // 3
    for i, (name, desc) in enumerate(methods):
        x = CX + i * (m2_w + IN(0.15))
        add_shape(slide, MSO_SHAPE.RECTANGLE, x, m2_top, m2_w, IN(0.42),
                  fill=PRIMARY, line_color=None,
                  text=name, font_size=12, bold=True, color=WHITE)
        body = add_shape(slide, MSO_SHAPE.RECTANGLE, x, m2_top + IN(0.42), m2_w, m2_h - IN(0.42),
                         fill=GRAY_LIGHT, line_color=GRAY_BORD)
        tf = body.text_frame; tf.word_wrap = True
        tf.margin_left = IN(0.1); tf.margin_right = IN(0.1)
        tf.margin_top = IN(0.06); tf.margin_bottom = IN(0.06)
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        r = p.add_run(); r.text = desc
        r.font.size = Pt(9.5); r.font.color.rgb = GRAY_DARK; r.font.name = "Calibri"


def fill_slide_metrics_formulas(slide):
    """Slide 7 [НОВЕ]: Метрики оцінювання — формули."""
    set_title(slide, "Метрики оцінювання retrieval-якості")
    remove_body_placeholder(slide)

    # 5 metric formula cards
    formulas = [
        ("nDCG@k",
         "DCG@k  /  IDCG@k",
         "Нормалізована якість ранжування",
         "де DCG@k = Σᵢ relᵢ / log₂(i+1),\nIDCG@k — DCG ідеального ранжування",
         PRIMARY),
        ("MRR",
         "(1/|Q|) · Σ (1 / rankq)",
         "Mean Reciprocal Rank",
         "середнє по запитах: 1 / (позиція\nпершого релевантного документа)",
         PRIMARY_LT),
        ("Recall@k",
         "|Rel ∩ Retrieved|  /  |Rel|",
         "Повнота пошуку",
         "частка знайдених релевантних\nдокументів серед усіх релевантних",
         TEAL),
        ("P@k",
         "|Relk|  /  k",
         "Precision@k",
         "частка релевантних серед\nперших k результатів",
         PURPLE),
    ]
    card_w = (CW - IN(0.3)) // 4
    card_h = CH - IN(0.1)
    for i, (name, formula, label, desc, color) in enumerate(formulas):
        x = CX + i * (card_w + IN(0.1))
        # Header
        add_shape(slide, MSO_SHAPE.RECTANGLE, x, CY, card_w, IN(0.45),
                  fill=color, line_color=None,
                  text=name, font_size=15, bold=True, color=WHITE)
        # Body
        body_top = CY + IN(0.45)
        body_h = card_h - IN(0.45)
        add_shape(slide, MSO_SHAPE.RECTANGLE, x, body_top, card_w, body_h,
                  fill=GRAY_LIGHT, line_color=GRAY_BORD)
        # Formula highlight box (inside body)
        formula_box = add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE,
                                x + IN(0.15), body_top + IN(0.18),
                                card_w - IN(0.3), IN(0.7),
                                fill=WHITE, line_color=color, line_width=Pt(1.5))
        ftf = formula_box.text_frame
        ftf.word_wrap = True
        ftf.margin_left = IN(0.06); ftf.margin_right = IN(0.06)
        ftf.margin_top = IN(0.04); ftf.margin_bottom = IN(0.04)
        ftf.vertical_anchor = MSO_ANCHOR.MIDDLE
        fp = ftf.paragraphs[0]; fp.alignment = PP_ALIGN.CENTER
        fr = fp.add_run(); fr.text = formula
        fr.font.size = Pt(12); fr.font.bold = True; fr.font.color.rgb = color
        fr.font.name = "Cambria Math"
        # Label
        add_text(slide, label,
                 x + IN(0.1), body_top + IN(1.05),
                 card_w - IN(0.2), IN(0.32),
                 font_size=11, bold=True, color=GRAY_DARK, align=PP_ALIGN.CENTER)
        # Description
        add_text(slide, desc,
                 x + IN(0.1), body_top + IN(1.4),
                 card_w - IN(0.2), IN(0.7),
                 font_size=9.5, color=GRAY_MED, align=PP_ALIGN.CENTER)


def fill_slide_mcda(slide):
    """Slide 8 [НОВЕ]: МКВ — методологія."""
    set_title(slide, "Багатокритеріальний вибір моделі (МКВ)")
    remove_body_placeholder(slide)

    # 3 method cards on top
    methods = [
        ("1. Парето-домінування",
         "P = { aᵢ ∈ A | ∄ ak ∈ A : ak ≻ aᵢ }",
         "Виявлення альтернатив, які не домінуються жодною іншою за усіма критеріями одночасно"),
        ("2. Нормалізація шкал",
         "yᵢⱼ = (xᵢⱼ − min xⱼ) / (max xⱼ − min xⱼ)",
         "Приведення усіх критеріїв до єдиної шкали [0; 1] з врахуванням напряму оптимізації"),
        ("3. Лінійна адитивна згортка",
         "U(aᵢ) = Σⱼ wⱼ · yᵢⱼ ,   де Σⱼ wⱼ = 1",
         "Інтегральна оцінка корисності з ваговими коефіцієнтами критеріїв"),
    ]
    method_h = IN(0.85)
    method_w = (CW - IN(0.3)) // 3
    for i, (name, formula, desc) in enumerate(methods):
        x = CX + i * (method_w + IN(0.15))
        # Header
        add_shape(slide, MSO_SHAPE.RECTANGLE, x, CY, method_w, IN(0.4),
                  fill=PRIMARY, line_color=None,
                  text=name, font_size=11, bold=True, color=WHITE)
        # Body
        body = add_shape(slide, MSO_SHAPE.RECTANGLE, x, CY + IN(0.4), method_w, method_h - IN(0.4),
                         fill=GRAY_LIGHT, line_color=GRAY_BORD)
        tf = body.text_frame; tf.word_wrap = True
        tf.margin_left = IN(0.1); tf.margin_right = IN(0.1)
        tf.margin_top = IN(0.06); tf.margin_bottom = IN(0.06)
        tf.vertical_anchor = MSO_ANCHOR.TOP
        # Formula
        p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        r = p.add_run(); r.text = formula
        r.font.size = Pt(11); r.font.bold = True; r.font.color.rgb = ACCENT; r.font.name = "Cambria Math"
        # Description
        p = tf.add_paragraph(); p.alignment = PP_ALIGN.CENTER
        p.space_before = Pt(4)
        r = p.add_run(); r.text = desc
        r.font.size = Pt(9); r.font.color.rgb = GRAY_DARK; r.font.name = "Calibri"

    # Weight assignment table
    table_y = CY + method_h + IN(0.2)
    add_text(slide, "ДОМЕННО-СПЕЦИФІЧНІ ПРОФІЛІ ВАГ КРИТЕРІЇВ", CX, table_y, CW, IN(0.25),
             font_size=11, bold=True, color=PRIMARY, align=PP_ALIGN.LEFT)
    # Build a small table
    rows = 4; cols = 6
    table_top = table_y + IN(0.28)
    table_h = IN(1.3)
    table_shape = slide.shapes.add_table(rows, cols, CX, table_top, CW, table_h)
    tbl = table_shape.table

    headers = ["Домен", "nDCG@k", "MRR@k", "Recall@k", "P@k", "Latency (ms)"]
    weight_data = [
        ("Технічний", "0.30", "0.20", "0.25", "0.05", "0.20"),
        ("Юридичний", "0.30", "0.30", "0.20", "0.10", "0.10"),
        ("Медичний", "0.25", "0.15", "0.35", "0.10", "0.15"),
    ]
    # Headers
    for j, h in enumerate(headers):
        cell = tbl.cell(0, j)
        cell.fill.solid(); cell.fill.fore_color.rgb = PRIMARY
        cell.text = ""
        p = cell.text_frame.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        r = p.add_run(); r.text = h
        r.font.size = Pt(11); r.font.bold = True; r.font.color.rgb = WHITE; r.font.name = "Calibri"
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE

    domain_colors = [DOMAIN_TECH, DOMAIN_LEGAL, DOMAIN_MED]
    for i, row_data in enumerate(weight_data):
        for j, val in enumerate(row_data):
            cell = tbl.cell(i + 1, j)
            cell.fill.solid()
            if j == 0:
                cell.fill.fore_color.rgb = domain_colors[i]
                col = WHITE
                bold = True
            else:
                cell.fill.fore_color.rgb = WHITE if i % 2 == 0 else GRAY_LIGHT
                col = GRAY_DARK
                bold = False
            cell.text = ""
            p = cell.text_frame.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
            r = p.add_run(); r.text = val
            r.font.size = Pt(11); r.font.bold = bold; r.font.color.rgb = col; r.font.name = "Calibri"
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE


def fill_slide_architecture(slide):
    """Slide 9: Архітектура системи."""
    set_title(slide, "Архітектура системи семантичного пошуку")
    remove_body_placeholder(slide)

    pipe_y = CY + IN(0.5)
    pipe_h = IN(0.85)
    boxes_data = [
        ("Документи", "PDF, DOCX,\nTXT, MD", PRIMARY),
        ("Чанкінг", "~600 слів,\nперекриття", PRIMARY_LT),
        ("Ембеддінг", "BGE-M3, E5,\nnomic, Qwen3", TEAL),
        ("FAISS\nіндекс", "IndexFlatIP,\nL2-нормалізація", PURPLE),
        ("Top-K", "Cosine\nsimilarity", ACCENT),
    ]
    total_box_w = IN(1.55)
    gap_w = IN(0.18)
    total_pw = total_box_w * 5 + gap_w * 4
    start_x = CX + (CW - total_pw) // 2
    for i, (name, desc, color) in enumerate(boxes_data):
        x = start_x + i * (total_box_w + gap_w)
        add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, x, pipe_y, total_box_w, IN(0.4),
                  fill=color, line_color=None,
                  text=name, font_size=12, bold=True, color=WHITE)
        body = add_shape(slide, MSO_SHAPE.RECTANGLE, x, pipe_y + IN(0.4), total_box_w, pipe_h - IN(0.4),
                         fill=GRAY_LIGHT, line_color=GRAY_BORD)
        tf = body.text_frame; tf.word_wrap = True
        tf.margin_left = IN(0.08); tf.margin_right = IN(0.08)
        tf.margin_top = IN(0.06); tf.margin_bottom = IN(0.06)
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        r = p.add_run(); r.text = desc
        r.font.size = Pt(9); r.font.color.rgb = GRAY_DARK; r.font.name = "Calibri"
        if i < len(boxes_data) - 1:
            arrow_y = pipe_y + pipe_h // 2
            add_arrow(slide, x + total_box_w, arrow_y,
                      x + total_box_w + gap_w, arrow_y, color=GRAY_DARK, weight=2.5)

    q_y = pipe_y + pipe_h + IN(0.35)
    add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, start_x, q_y, total_box_w, IN(0.55),
              fill=ACCENT, line_color=None,
              text="Запит\nкористувача", font_size=11, bold=True, color=WHITE)
    emb_x = start_x + 2 * (total_box_w + gap_w) + total_box_w // 2
    add_arrow(slide, start_x + total_box_w // 2, q_y, emb_x, pipe_y + pipe_h, color=ACCENT, weight=2.5)

    flask_x = start_x + 4 * (total_box_w + gap_w)
    add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, flask_x, q_y, total_box_w, IN(0.55),
              fill=PRIMARY, line_color=None,
              text="Flask Web UI", font_size=11, bold=True, color=WHITE)
    add_arrow(slide, flask_x + total_box_w // 2, pipe_y + pipe_h,
              flask_x + total_box_w // 2, q_y, color=PRIMARY, weight=2.5)

    note_y = q_y + IN(0.7)
    add_text(slide,
             "Основні компоненти:  build_index.py  •  evaluate_benchmark.py  •  embedding_models.py  •  run_all_benchmarks.py  •  app.py",
             CX, note_y, CW, IN(0.3),
             font_size=9.5, italic=True, color=GRAY_MED, align=PP_ALIGN.CENTER)


def fill_slide_software(slide):
    """Slide 10: Програмне забезпечення."""
    set_title(slide, "Програмне забезпечення та технології")
    remove_body_placeholder(slide)

    stacks = [
        ("ML / EMBEDDING", PRIMARY, [
            "Python 3.11", "sentence-transformers", "transformers (HF)",
            "PyTorch", "trust_remote_code (nomic, Qwen3)", "Cambria Math (формули)",
        ]),
        ("RETRIEVAL / DATA", TEAL, [
            "FAISS (IndexFlatIP)", "rank_bm25 (Okapi)", "L2-нормалізація",
            "Чанкінг (~600 слів)", "PDF / DOCX / TXT / MD",
            "JSONL (chunks, qrels, queries)",
        ]),
        ("WEB / EVAL / ANALYSIS", ACCENT, [
            "Flask (Python)", "NumPy, scikit-learn",
            "Bootstrap CI (n=2000)", "Pareto + лінійна згортка",
            "Власний benchmark harness", "Git + GitHub",
        ]),
    ]
    col_w = (CW - IN(0.3)) // 3
    col_h = CH - IN(0.1)
    col_top = CY
    for i, (name, color, items) in enumerate(stacks):
        x = CX + i * (col_w + IN(0.15))
        add_shape(slide, MSO_SHAPE.RECTANGLE, x, col_top, col_w, IN(0.5),
                  fill=color, line_color=None,
                  text=name, font_size=14, bold=True, color=WHITE)
        body = add_shape(slide, MSO_SHAPE.RECTANGLE, x, col_top + IN(0.5), col_w, col_h - IN(0.5),
                         fill=GRAY_LIGHT, line_color=GRAY_BORD)
        tf = body.text_frame; tf.word_wrap = True
        tf.margin_left = IN(0.18); tf.margin_right = IN(0.1)
        tf.margin_top = IN(0.12); tf.margin_bottom = IN(0.1)
        tf.vertical_anchor = MSO_ANCHOR.TOP
        p = tf.paragraphs[0]
        for j, item in enumerate(items):
            if j == 0:
                para = p
            else:
                para = tf.add_paragraph()
            para.alignment = PP_ALIGN.LEFT
            para.space_after = Pt(4)
            r = para.add_run(); r.text = "▸ " + item
            r.font.size = Pt(11); r.font.color.rgb = GRAY_DARK; r.font.name = "Calibri"


def fill_slide_dataset(slide):
    """Slide 11: Benchmark dataset."""
    set_title(slide, "Benchmark dataset: україномовні доменні колекції")
    remove_body_placeholder(slide)

    add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, CX, CY, CW, IN(0.5),
              fill=PRIMARY, line_color=None,
              text="Власноруч сформований domain-specific україномовний benchmark dataset",
              font_size=12, bold=True, color=WHITE)

    domains = [
        ("Технічний", DOMAIN_TECH, "10", "178", "100",
         "програмна інженерія, ML, RAG,\nаномалії, 3D-друк, лідарні системи"),
        ("Юридичний", DOMAIN_LEGAL, "13", "1 046", "100",
         "кодекси, нормативно-правові\nакти, законодавство"),
        ("Медичний", DOMAIN_MED, "9", "267", "100",
         "клінічні протоколи,\nмедичні рекомендації"),
    ]
    d_top = CY + IN(0.65)
    d_h = IN(2.0)
    d_w = (CW - IN(0.3)) // 3
    for i, (name, color, docs, chunks, queries, topics) in enumerate(domains):
        x = CX + i * (d_w + IN(0.15))
        add_shape(slide, MSO_SHAPE.RECTANGLE, x, d_top, d_w, IN(0.4),
                  fill=color, line_color=None,
                  text=name, font_size=14, bold=True, color=WHITE)
        body = add_shape(slide, MSO_SHAPE.RECTANGLE, x, d_top + IN(0.4), d_w, d_h - IN(0.4),
                         fill=GRAY_LIGHT, line_color=GRAY_BORD)
        tf = body.text_frame; tf.word_wrap = True
        tf.margin_left = IN(0.12); tf.margin_right = IN(0.12)
        tf.margin_top = IN(0.1); tf.margin_bottom = IN(0.1)
        tf.vertical_anchor = MSO_ANCHOR.TOP
        p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        r = p.add_run(); r.text = docs + "  /  " + chunks + "  /  " + queries
        r.font.size = Pt(20); r.font.bold = True; r.font.color.rgb = color; r.font.name = "Calibri"
        p = tf.add_paragraph(); p.alignment = PP_ALIGN.CENTER
        p.space_after = Pt(8)
        r = p.add_run(); r.text = "документи  /  чанки  /  запити"
        r.font.size = Pt(8); r.font.italic = True; r.font.color.rgb = GRAY_MED; r.font.name = "Calibri"
        p = tf.add_paragraph(); p.alignment = PP_ALIGN.LEFT
        r = p.add_run(); r.text = "Тематика:"
        r.font.size = Pt(9); r.font.bold = True; r.font.color.rgb = GRAY_DARK; r.font.name = "Calibri"
        p = tf.add_paragraph(); p.alignment = PP_ALIGN.LEFT
        r = p.add_run(); r.text = topics
        r.font.size = Pt(9); r.font.color.rgb = GRAY_DARK; r.font.name = "Calibri"

    tot_y = d_top + d_h + IN(0.15)
    tot_box = add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, CX, tot_y, CW, IN(0.5),
                        fill=ACCENT, line_color=None)
    tf = tot_box.text_frame; tf.word_wrap = True
    tf.margin_left = IN(0.2); tf.margin_right = IN(0.2)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = "ВСЬОГО: "
    r.font.size = Pt(11); r.font.bold = True; r.font.color.rgb = WHITE; r.font.name = "Calibri"
    r = p.add_run(); r.text = "32 документи · 1 491 чанк · 300 запитів · qrels сформовані вручну на рівні чанків"
    r.font.size = Pt(11); r.font.color.rgb = WHITE; r.font.name = "Calibri"


def fill_slide_sample_queries(slide):
    """Slide 12 [НОВЕ]: Приклади benchmark-запитів."""
    set_title(slide, "Приклади benchmark-запитів")
    remove_body_placeholder(slide)

    # 3 columns per domain
    queries = [
        ("Технічний", DOMAIN_TECH, [
            ("q1", "виявлення аномалій у фінансових транзакціях"),
            ("q3", "комбінація Autoencoder та Isolation Forest для fraud detection"),
            ("q6", "прогноз відтоку клієнтів (churn prediction)"),
            ("q7", "когортний аналіз для сегментації клієнтів"),
            ("q8", "етичні питання використання персональних даних"),
        ]),
        ("Юридичний", DOMAIN_LEGAL, [
            ("L001", "що таке цивільна дієздатність фізичної особи"),
            ("L002", "з якого віку настає повна цивільна дієздатність"),
            ("L003", "підстави обмеження цивільної дієздатності судом"),
            ("L004", "поняття юридичної особи та порядок її реєстрації"),
            ("L005", "що входить до складу спадщини"),
        ]),
        ("Медичний", DOMAIN_MED, [
            ("M001", "що таке серцево-судинна система"),
            ("M002", "будова та функції серця людини"),
            ("M003", "що таке артеріальний тиск і як він регулюється"),
            ("M004", "велике і мале коло кровообігу — відмінності"),
            ("M005", "будова дихальної системи людини"),
        ]),
    ]
    col_w = (CW - IN(0.3)) // 3
    col_h = IN(2.9)
    for i, (name, color, items) in enumerate(queries):
        x = CX + i * (col_w + IN(0.15))
        # Header
        add_shape(slide, MSO_SHAPE.RECTANGLE, x, CY, col_w, IN(0.4),
                  fill=color, line_color=None,
                  text=name, font_size=13, bold=True, color=WHITE)
        # Body
        body = add_shape(slide, MSO_SHAPE.RECTANGLE, x, CY + IN(0.4), col_w, col_h - IN(0.4),
                         fill=GRAY_LIGHT, line_color=GRAY_BORD)
        tf = body.text_frame; tf.word_wrap = True
        tf.margin_left = IN(0.15); tf.margin_right = IN(0.12)
        tf.margin_top = IN(0.1); tf.margin_bottom = IN(0.1)
        tf.vertical_anchor = MSO_ANCHOR.TOP
        p = tf.paragraphs[0]
        for j, (qid, qtext) in enumerate(items):
            if j == 0:
                para = p
            else:
                para = tf.add_paragraph()
            para.alignment = PP_ALIGN.LEFT
            para.space_after = Pt(7)
            r = para.add_run(); r.text = qid + " "
            r.font.size = Pt(8); r.font.bold = True; r.font.color.rgb = color; r.font.name = "Consolas"
            r = para.add_run(); r.text = "« " + qtext + " »"
            r.font.size = Pt(9.5); r.font.color.rgb = GRAY_DARK; r.font.name = "Calibri"

    # Bottom note
    note_y = CY + col_h + IN(0.15)
    note_box = add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, CX, note_y, CW, IN(0.45),
                         fill=GRAY_LIGHT, line_color=ACCENT, line_width=Pt(1.5))
    tf = note_box.text_frame; tf.word_wrap = True
    tf.margin_left = IN(0.15); tf.margin_right = IN(0.15)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.LEFT
    r = p.add_run(); r.text = "Типи запитів: "
    r.font.size = Pt(10); r.font.bold = True; r.font.color.rgb = ACCENT; r.font.name = "Calibri"
    r = p.add_run()
    r.text = ("definition · factual · procedural · technical · topic · method · policy · comparison    "
              "(qrels — релевантні чанки сформовано вручну для кожного запиту)")
    r.font.size = Pt(9.5); r.font.color.rgb = GRAY_DARK; r.font.name = "Calibri"


def fill_slide_results(slide):
    """Slide 12: Результати — nDCG@10."""
    set_title(slide, "Результати експерименту: nDCG@10")
    remove_body_placeholder(slide)

    add_text(slide, "Зведені значення nDCG@10 (вище — краще; найкращий результат у домені виділено)",
             CX, CY, CW, IN(0.3),
             font_size=11, italic=True, color=GRAY_MED, align=PP_ALIGN.LEFT)

    table_y = CY + IN(0.3)
    table_h = IN(2.4)
    rows = 6
    cols = 5
    table_shape = slide.shapes.add_table(rows, cols, CX, table_y, CW, table_h)
    tbl = table_shape.table

    tbl.columns[0].width = IN(3.0)
    for i in range(1, 5):
        tbl.columns[i].width = (CW - IN(3.0)) // 4

    headers = ["Модель", "Технічний", "Юридичний", "Медичний", "Середнє"]
    for j, h in enumerate(headers):
        cell = tbl.cell(0, j)
        cell.fill.solid(); cell.fill.fore_color.rgb = PRIMARY
        cell.text = ""
        p = cell.text_frame.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        r = p.add_run(); r.text = h
        r.font.size = Pt(13); r.font.bold = True; r.font.color.rgb = WHITE; r.font.name = "Calibri"
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE

    data = [
        ("BGE-M3 (BAAI)",            0.6722, 0.3065, 0.4339),
        ("Qwen3-Embedding-0.6B",     0.6325, 0.3199, 0.3629),
        ("E5-base (multilingual)",   0.6121, 0.2567, 0.3909),
        ("BM25 (Okapi) — базелайн",  0.4861, 0.1875, 0.3222),
        ("nomic-embed-text-v1.5",    0.3765, 0.0951, 0.1668),
    ]
    tech_vals = [d[1] for d in data]; max_tech = max(tech_vals)
    leg_vals = [d[2] for d in data]; max_leg = max(leg_vals)
    med_vals = [d[3] for d in data]; max_med = max(med_vals)

    for i, (model, tech, leg, med) in enumerate(data):
        avg = (tech + leg + med) / 3
        bg = WHITE if i % 2 == 0 else GRAY_LIGHT
        row_data = [
            (model, False, GRAY_DARK, PP_ALIGN.LEFT, bg),
            (f"{tech:.4f}", tech == max_tech, SUCCESS if tech == max_tech else GRAY_DARK,
             PP_ALIGN.CENTER, SUCCESS if tech == max_tech else bg),
            (f"{leg:.4f}",  leg == max_leg, SUCCESS if leg == max_leg else GRAY_DARK,
             PP_ALIGN.CENTER, SUCCESS if leg == max_leg else bg),
            (f"{med:.4f}",  med == max_med, SUCCESS if med == max_med else GRAY_DARK,
             PP_ALIGN.CENTER, SUCCESS if med == max_med else bg),
            (f"{avg:.4f}", False, GRAY_DARK, PP_ALIGN.CENTER, bg),
        ]
        for j, (txt, is_best, fg, align, cell_bg) in enumerate(row_data):
            cell = tbl.cell(i + 1, j)
            cell.fill.solid(); cell.fill.fore_color.rgb = cell_bg
            cell.text = ""
            p = cell.text_frame.paragraphs[0]; p.alignment = align
            r = p.add_run(); r.text = txt
            r.font.size = Pt(11)
            r.font.bold = is_best or (j == 0)
            r.font.color.rgb = WHITE if (is_best and cell_bg == SUCCESS) else fg
            r.font.name = "Calibri"
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE

    ci_y = table_y + table_h + IN(0.1)
    ci_box = add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, CX, ci_y, CW, IN(0.55),
                       fill=GRAY_LIGHT, line_color=ACCENT, line_width=Pt(1.5))
    tf = ci_box.text_frame; tf.word_wrap = True
    tf.margin_left = IN(0.15); tf.margin_right = IN(0.15)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.LEFT
    r = p.add_run(); r.text = "Bootstrap 95% CI (n = 2 000): "
    r.font.size = Pt(10); r.font.bold = True; r.font.color.rgb = ACCENT; r.font.name = "Calibri"
    r = p.add_run()
    r.text = ("BGE-M3 [0.606 ; 0.737]  vs  BM25 [0.431 ; 0.541]  →  інтервали не перетинаються "
              "→ статистично значуща перевага семантичного пошуку над лексичним")
    r.font.size = Pt(10); r.font.color.rgb = GRAY_DARK; r.font.name = "Calibri"


def fill_slide_detailed_metrics(slide):
    """Slide 13 [НОВЕ]: Деталізовані метрики (MRR, Recall, P, ms) по доменах."""
    set_title(slide, "Деталізовані метрики за доменами")
    remove_body_placeholder(slide)

    # Three bar charts side by side: MRR, Recall, P (excluding nDCG which is on slide 12)
    chart_w = (CW - IN(0.4)) // 3
    chart_h = IN(2.4)
    chart_y = CY + IN(0.2)

    # Чарт MRR
    chart_data = CategoryChartData()
    chart_data.categories = ["Тех.", "Юр.", "Мед."]
    chart_data.add_series("BGE-M3",  (0.6993, 0.3831, 0.5054))
    chart_data.add_series("Qwen3",   (0.6640, 0.3999, 0.4184))
    chart_data.add_series("E5-base", (0.6340, 0.3145, 0.4888))
    chart_data.add_series("BM25",    (0.4808, 0.2533, 0.4003))
    chart_data.add_series("nomic",   (0.4122, 0.1384, 0.2110))
    add_chart_helper(slide, chart_data, "MRR@10", CX, chart_y, chart_w, chart_h)

    # Чарт Recall
    chart_data = CategoryChartData()
    chart_data.categories = ["Тех.", "Юр.", "Мед."]
    chart_data.add_series("BGE-M3",  (0.815, 0.375, 0.545))
    chart_data.add_series("Qwen3",   (0.7792, 0.375, 0.4483))
    chart_data.add_series("E5-base", (0.78, 0.35, 0.4583))
    chart_data.add_series("BM25",    (0.7542, 0.2333, 0.3983))
    chart_data.add_series("nomic",   (0.4725, 0.125, 0.1917))
    add_chart_helper(slide, chart_data, "Recall@10", CX + chart_w + IN(0.2), chart_y, chart_w, chart_h)

    # Чарт P
    chart_data = CategoryChartData()
    chart_data.categories = ["Тех.", "Юр.", "Мед."]
    chart_data.add_series("BGE-M3",  (0.143, 0.106, 0.140))
    chart_data.add_series("Qwen3",   (0.136, 0.107, 0.115))
    chart_data.add_series("E5-base", (0.136, 0.099, 0.114))
    chart_data.add_series("BM25",    (0.134, 0.064, 0.100))
    chart_data.add_series("nomic",   (0.075, 0.037, 0.051))
    add_chart_helper(slide, chart_data, "P@10", CX + 2 * (chart_w + IN(0.2)), chart_y, chart_w, chart_h)

    # Bottom: latency note
    note_y = chart_y + chart_h + IN(0.15)
    note_box = add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, CX, note_y, CW, IN(0.55),
                         fill=GRAY_LIGHT, line_color=ACCENT, line_width=Pt(1.5))
    tf = note_box.text_frame; tf.word_wrap = True
    tf.margin_left = IN(0.15); tf.margin_right = IN(0.15)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.LEFT
    r = p.add_run(); r.text = "Швидкодія (ms/query): "
    r.font.size = Pt(10); r.font.bold = True; r.font.color.rgb = ACCENT; r.font.name = "Calibri"
    r = p.add_run()
    r.text = ("BM25 ≈ 0.6–2.9   ◇   E5-base ≈ 67–225   ◇   nomic ≈ 131–407   "
              "◇   Qwen3 ≈ 365–466   ◇   BGE-M3 ≈ 203–2960")
    r.font.size = Pt(10); r.font.color.rgb = GRAY_DARK; r.font.name = "Calibri"


def add_chart_helper(slide, chart_data, title, x, y, w, h):
    """Створює один bar chart."""
    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, w, h, chart_data
    ).chart
    chart.has_title = True
    chart.chart_title.text_frame.text = title
    for run in chart.chart_title.text_frame.paragraphs[0].runs:
        run.font.size = Pt(11); run.font.bold = True; run.font.color.rgb = PRIMARY
    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    chart.legend.include_in_layout = False
    chart.legend.font.size = Pt(8)
    series_colors = [SUCCESS, ACCENT, PRIMARY_LT, GRAY_MED, DANGER]
    for i, series in enumerate(chart.plots[0].series):
        fill = series.format.fill
        fill.solid()
        fill.fore_color.rgb = series_colors[i % len(series_colors)]


def fill_slide_analysis(slide):
    """Slide 14: МКВ + bar chart of U-scores."""
    set_title(slide, "Багатокритеріальний вибір моделі (МКВ): результати")
    remove_body_placeholder(slide)

    lw = IN(3.5)
    pareto_box = add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, CX, CY, lw, IN(1.1),
                           fill=PRIMARY, line_color=None)
    tf = pareto_box.text_frame; tf.word_wrap = True
    tf.margin_left = IN(0.15); tf.margin_right = IN(0.15)
    tf.margin_top = IN(0.1); tf.margin_bottom = IN(0.1)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.LEFT
    r = p.add_run(); r.text = "ПАРЕТО-ОПТИМАЛЬНІ"
    r.font.size = Pt(11); r.font.bold = True; r.font.color.rgb = ACCENT; r.font.name = "Calibri"
    p = tf.add_paragraph(); p.alignment = PP_ALIGN.LEFT
    r = p.add_run(); r.text = "P = { BGE-M3,  E5-base }"
    r.font.size = Pt(15); r.font.bold = True; r.font.color.rgb = WHITE; r.font.name = "Calibri"

    findings = [
        ("BGE-M3", "лідер у всіх 3 доменах за nDCG / MRR / Recall", SUCCESS),
        ("E5-base", "швидка альтернатива (~225 мс vs ~2960 мс у BGE-M3)", PRIMARY_LT),
        ("Qwen3", "конкурентна якість, але ~1500 мс на CPU — потрібен GPU", ACCENT),
        ("nomic", "найслабші результати; в окремих доменах поступається BM25", DANGER),
    ]
    f_top = CY + IN(1.25)
    f_h = IN(0.5); f_gap = IN(0.05)
    for i, (name, desc, color) in enumerate(findings):
        y = f_top + i * (f_h + f_gap)
        add_shape(slide, MSO_SHAPE.RECTANGLE, CX, y, IN(0.12), f_h,
                  fill=color, line_color=None)
        box = add_shape(slide, MSO_SHAPE.RECTANGLE, CX + IN(0.12), y, lw - IN(0.12), f_h,
                        fill=GRAY_LIGHT, line_color=GRAY_BORD)
        tf = box.text_frame; tf.word_wrap = True
        tf.margin_left = IN(0.12); tf.margin_right = IN(0.1)
        tf.margin_top = IN(0.06); tf.margin_bottom = IN(0.06)
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]; p.alignment = PP_ALIGN.LEFT
        r = p.add_run(); r.text = name + ":  "
        r.font.size = Pt(11); r.font.bold = True; r.font.color.rgb = color; r.font.name = "Calibri"
        r = p.add_run(); r.text = desc
        r.font.size = Pt(10); r.font.color.rgb = GRAY_DARK; r.font.name = "Calibri"

    chart_x = CX + lw + IN(0.2)
    chart_w = CW - lw - IN(0.2)

    chart_data = CategoryChartData()
    chart_data.categories = ["Технічний", "Юридичний", "Медичний"]
    chart_data.add_series("BGE-M3", (0.920, 0.927, 0.901))
    chart_data.add_series("E5-base", (0.842, 0.770, 0.745))
    chart_data.add_series("Qwen3", (0.650, 0.770, 0.580))
    chart_data.add_series("BM25", (0.425, 0.350, 0.400))
    chart_data.add_series("nomic", (0.350, 0.180, 0.220))

    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED, chart_x, CY, chart_w, CH, chart_data
    ).chart
    chart.has_title = True
    chart.chart_title.text_frame.text = "Інтегральна оцінка U(aᵢ) за доменами"
    for run in chart.chart_title.text_frame.paragraphs[0].runs:
        run.font.size = Pt(11); run.font.bold = True; run.font.color.rgb = PRIMARY
    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    chart.legend.include_in_layout = False
    chart.legend.font.size = Pt(10)
    series_colors = [SUCCESS, PRIMARY_LT, ACCENT, GRAY_MED, DANGER]
    for i, series in enumerate(chart.plots[0].series):
        fill = series.format.fill
        fill.solid()
        fill.fore_color.rgb = series_colors[i]


def fill_slide_quality_speed(slide):
    """Slide 15 [НОВЕ]: Якість vs швидкодія — scatter."""
    set_title(slide, "Аналіз: якість пошуку vs швидкодія")
    remove_body_placeholder(slide)

    # Bubble chart: x = log10(ms), y = nDCG@10
    # 5 моделей × 3 домени = 15 точок
    # Краще зробити XY scatter chart

    # Build XY scatter data
    chart_data = XyChartData()

    # Series per model (5 series)
    bge_series = chart_data.add_series("BGE-M3")
    bge_series.add_data_point(2960, 0.6722)
    bge_series.add_data_point(202.9, 0.3065)
    bge_series.add_data_point(236.3, 0.4339)

    e5_series = chart_data.add_series("E5-base")
    e5_series.add_data_point(225.1, 0.6121)
    e5_series.add_data_point(67.0, 0.2567)
    e5_series.add_data_point(80.3, 0.3909)

    qwen_series = chart_data.add_series("Qwen3")
    qwen_series.add_data_point(365.4, 0.6325)
    qwen_series.add_data_point(399.1, 0.3199)
    qwen_series.add_data_point(465.9, 0.3629)

    bm25_series = chart_data.add_series("BM25")
    bm25_series.add_data_point(1.5, 0.4861)
    bm25_series.add_data_point(2.9, 0.1875)
    bm25_series.add_data_point(0.6, 0.3222)

    nomic_series = chart_data.add_series("nomic")
    nomic_series.add_data_point(407.4, 0.3765)
    nomic_series.add_data_point(131.6, 0.0951)
    nomic_series.add_data_point(150.4, 0.1668)

    # Chart on left, takeaways on right
    chart_w = IN(6.0)
    chart_h = CH
    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.XY_SCATTER, CX, CY, chart_w, chart_h, chart_data
    ).chart
    chart.has_title = False
    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    chart.legend.include_in_layout = False
    chart.legend.font.size = Pt(10)

    # Axis labels
    cat_axis = chart.category_axis
    cat_axis.has_title = True
    cat_axis.axis_title.text_frame.text = "Час відповіді (мс/запит, log scale)"
    for r in cat_axis.axis_title.text_frame.paragraphs[0].runs:
        r.font.size = Pt(10); r.font.color.rgb = GRAY_DARK
    # Set log scale on category axis (which is X for XY scatter)
    cat_axis_xml = cat_axis._element
    scaling = cat_axis_xml.find(qn("c:scaling"))
    if scaling is not None:
        log_base = scaling.find(qn("c:logBase"))
        if log_base is None:
            log_base = etree.SubElement(scaling, qn("c:logBase"))
            log_base.set("val", "10")

    val_axis = chart.value_axis
    val_axis.has_title = True
    val_axis.axis_title.text_frame.text = "nDCG@10"
    for r in val_axis.axis_title.text_frame.paragraphs[0].runs:
        r.font.size = Pt(10); r.font.color.rgb = GRAY_DARK

    # Color points
    series_colors = [SUCCESS, PRIMARY_LT, ACCENT, GRAY_MED, DANGER]
    for i, series in enumerate(chart.plots[0].series):
        marker = series.marker
        marker.style = 8  # CIRCLE
        marker.size = 12
        fill = marker.format.fill
        fill.solid()
        fill.fore_color.rgb = series_colors[i]
        line = marker.format.line
        line.color.rgb = WHITE
        # Hide the connecting line (it's XY scatter without line)
        ser_line = series.format.line
        ser_line.fill.background()

    # Right side: takeaways
    tk_x = CX + chart_w + IN(0.2)
    tk_w = CW - chart_w - IN(0.2)
    add_text(slide, "КЛЮЧОВІ ВИСНОВКИ", tk_x, CY, tk_w, IN(0.3),
             font_size=12, bold=True, color=PRIMARY, align=PP_ALIGN.LEFT)

    takeaways = [
        ("Лідер якості:", "BGE-M3", "найвища nDCG, але повільний", SUCCESS),
        ("Найшвидший:", "BM25", "<3 мс, але якість обмежена", GRAY_MED),
        ("Оптимум:", "E5-base", "Парето-frontier, ~225 мс", PRIMARY_LT),
        ("Без GPU:", "Qwen3", "повільний, конкурентний", ACCENT),
        ("Найслабший:", "nomic", "поступається навіть BM25", DANGER),
    ]
    tk_top = CY + IN(0.35)
    tk_h = IN(0.55); tk_gap = IN(0.06)
    for i, (label, name, desc, color) in enumerate(takeaways):
        y = tk_top + i * (tk_h + tk_gap)
        add_shape(slide, MSO_SHAPE.RECTANGLE, tk_x, y, IN(0.1), tk_h,
                  fill=color, line_color=None)
        box = add_shape(slide, MSO_SHAPE.RECTANGLE, tk_x + IN(0.1), y, tk_w - IN(0.1), tk_h,
                        fill=GRAY_LIGHT, line_color=GRAY_BORD)
        tf = box.text_frame; tf.word_wrap = True
        tf.margin_left = IN(0.1); tf.margin_right = IN(0.06)
        tf.margin_top = IN(0.04); tf.margin_bottom = IN(0.04)
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]; p.alignment = PP_ALIGN.LEFT
        r = p.add_run(); r.text = label + " "
        r.font.size = Pt(8.5); r.font.bold = True; r.font.color.rgb = GRAY_MED; r.font.name = "Calibri"
        r = p.add_run(); r.text = name
        r.font.size = Pt(10); r.font.bold = True; r.font.color.rgb = color; r.font.name = "Calibri"
        p = tf.add_paragraph(); p.alignment = PP_ALIGN.LEFT
        r = p.add_run(); r.text = desc
        r.font.size = Pt(8.5); r.font.color.rgb = GRAY_DARK; r.font.name = "Calibri"


def fill_slide_approbation(slide):
    """Slide 16: Апробація."""
    set_title(slide, "Апробація результатів дослідження")
    remove_body_placeholder(slide)

    confs = [
        ("Конференція №1   |   2026",
         "Innovative Research in Science and Economy",
         "2-я Міжнародна науково-практична конференція",
         "Стаття: «Порівняльний аналіз ефективності моделей векторних ембеддінгів для задач семантичного пошуку в корпоративних базах знань»",
         "представлено benchmark-методологію та порівняння 4 моделей у трьох доменних колекціях; обґрунтовано перевагу BGE-M3 над класичним BM25.",
         "у співавторстві з Русаковою Н.Є.",
         PRIMARY),
        ("Конференція №2   |   2026",
         "Радіоелектроніка та молодь у XXI столітті",
         "30-й Міжнародний молодіжний форум, ХНУРЕ",
         "Тези: «Підвищення якості семантичного пошуку в базах знань із використанням векторних подань»",
         "подано підхід до підвищення якості retrieval у корпоративних базах знань на основі контрастивно-навчених embedding-моделей.",
         "у співавторстві з Русаковою Н.Є.",
         TEAL),
    ]
    c_w = (CW - IN(0.3)) // 2
    c_h = IN(2.85)
    for i, (label, name, conf_type, paper, contrib, author, color) in enumerate(confs):
        x = CX + i * (c_w + IN(0.15))
        add_shape(slide, MSO_SHAPE.RECTANGLE, x, CY, c_w, IN(0.45),
                  fill=color, line_color=None,
                  text=label, font_size=12, bold=True, color=WHITE)
        body = add_shape(slide, MSO_SHAPE.RECTANGLE, x, CY + IN(0.45), c_w, c_h - IN(0.45),
                         fill=GRAY_LIGHT, line_color=GRAY_BORD)
        tf = body.text_frame; tf.word_wrap = True
        tf.margin_left = IN(0.18); tf.margin_right = IN(0.15)
        tf.margin_top = IN(0.15); tf.margin_bottom = IN(0.1)
        tf.vertical_anchor = MSO_ANCHOR.TOP
        p = tf.paragraphs[0]; p.alignment = PP_ALIGN.LEFT
        r = p.add_run(); r.text = "«" + name + "»"
        r.font.size = Pt(13); r.font.bold = True; r.font.color.rgb = color; r.font.name = "Calibri"
        p = tf.add_paragraph(); p.alignment = PP_ALIGN.LEFT; p.space_before = Pt(2); p.space_after = Pt(6)
        r = p.add_run(); r.text = conf_type
        r.font.size = Pt(9.5); r.font.italic = True; r.font.color.rgb = GRAY_MED; r.font.name = "Calibri"
        p = tf.add_paragraph(); p.alignment = PP_ALIGN.LEFT
        r = p.add_run(); r.text = "▸ " + paper
        r.font.size = Pt(10); r.font.color.rgb = GRAY_DARK; r.font.name = "Calibri"
        p = tf.add_paragraph(); p.alignment = PP_ALIGN.LEFT; p.space_before = Pt(6); p.space_after = Pt(6)
        r = p.add_run(); r.text = "Ключовий внесок: "
        r.font.size = Pt(9.5); r.font.bold = True; r.font.color.rgb = color; r.font.name = "Calibri"
        r = p.add_run(); r.text = contrib
        r.font.size = Pt(9.5); r.font.color.rgb = GRAY_DARK; r.font.name = "Calibri"
        p = tf.add_paragraph(); p.alignment = PP_ALIGN.LEFT; p.space_before = Pt(4)
        r = p.add_run(); r.text = author
        r.font.size = Pt(9); r.font.italic = True; r.font.color.rgb = GRAY_MED; r.font.name = "Calibri"

    # Bottom summary banner
    summary_y = CY + c_h + IN(0.15)
    sum_box = add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, CX, summary_y, CW, IN(0.55),
                        fill=ACCENT, line_color=None)
    tf = sum_box.text_frame; tf.word_wrap = True
    tf.margin_left = IN(0.2); tf.margin_right = IN(0.2)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = "ЗАГАЛОМ: "
    r.font.size = Pt(11); r.font.bold = True; r.font.color.rgb = WHITE; r.font.name = "Calibri"
    r = p.add_run(); r.text = "2 наукові праці  ·  2 виступи на міжнародних конференціях  ·  у співавторстві з науковим керівником  ·  ХНУРЕ, 2026"
    r.font.size = Pt(11); r.font.color.rgb = WHITE; r.font.name = "Calibri"


def fill_slide_summary(slide):
    """Slide 17: Підсумки."""
    set_title(slide, "Підсумки та практичні рекомендації")
    remove_body_placeholder(slide)

    hero = add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, CX, CY, CW, IN(0.55),
                     fill=PRIMARY, line_color=None)
    tf = hero.text_frame; tf.word_wrap = True
    tf.margin_left = IN(0.18); tf.margin_right = IN(0.18)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = "✓  ПОСТАВЛЕНУ МЕТУ ДОСЯГНУТО: "
    r.font.size = Pt(13); r.font.bold = True; r.font.color.rgb = ACCENT; r.font.name = "Calibri"
    r = p.add_run(); r.text = "обґрунтовано доцільність застосування embedding-моделей для семантичного пошуку"
    r.font.size = Pt(11); r.font.color.rgb = WHITE; r.font.name = "Calibri"

    recs = [
        ("BGE-M3", "основна модель для впровадження", "Найкраща якість, стабільність у різних доменах", SUCCESS),
        ("E5-base", "швидка альтернатива", "Прийнятна якість + ~13× швидше за BGE-M3", PRIMARY_LT),
        ("Qwen3", "лише з GPU", "Висока якість, але непрактична на CPU", ACCENT),
        ("BM25", "конкурентний baseline", "Залишається сильним у певних доменах", GRAY_MED),
    ]
    r_top = CY + IN(0.7)
    r_w = (CW - IN(0.15)) // 2
    r_h = IN(1.0)
    for i, (name, role, desc, color) in enumerate(recs):
        row = i // 2; col = i % 2
        x = CX + col * (r_w + IN(0.15))
        y = r_top + row * (r_h + IN(0.1))
        add_shape(slide, MSO_SHAPE.RECTANGLE, x, y, IN(0.18), r_h,
                  fill=color, line_color=None)
        body = add_shape(slide, MSO_SHAPE.RECTANGLE, x + IN(0.18), y, r_w - IN(0.18), r_h,
                         fill=GRAY_LIGHT, line_color=GRAY_BORD)
        tf = body.text_frame; tf.word_wrap = True
        tf.margin_left = IN(0.15); tf.margin_right = IN(0.12)
        tf.margin_top = IN(0.1); tf.margin_bottom = IN(0.1)
        tf.vertical_anchor = MSO_ANCHOR.TOP
        p = tf.paragraphs[0]; p.alignment = PP_ALIGN.LEFT
        r = p.add_run(); r.text = name
        r.font.size = Pt(15); r.font.bold = True; r.font.color.rgb = color; r.font.name = "Calibri"
        r = p.add_run(); r.text = "  —  " + role
        r.font.size = Pt(11); r.font.italic = True; r.font.color.rgb = GRAY_DARK; r.font.name = "Calibri"
        p = tf.add_paragraph(); p.alignment = PP_ALIGN.LEFT; p.space_before = Pt(4)
        r = p.add_run(); r.text = desc
        r.font.size = Pt(10); r.font.color.rgb = GRAY_DARK; r.font.name = "Calibri"

    thx_y = r_top + 2 * r_h + IN(0.2)
    add_text(slide, "Дякую за увагу!", CX, thx_y, CW, IN(0.45),
             font_size=22, bold=True, color=PRIMARY, align=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════════════════════
# MAIN
# ════════════════════════════════════════════════════════════════════════

def main():
    prs = Presentation(str(PPTX))
    print(f"Opened: {len(prs.slides)} slides")

    # 1. Дублюємо slide 2 (TITLE_AND_BODY) 6 разів — вони стають slide 13-18
    for _ in range(6):
        duplicate_slide(prs, 1)
    print(f"After duplication: {len(prs.slides)} slides")

    # 2. Реорганізовуємо порядок (0-indexed)
    # Original 0-11: original slides 1-12
    # New 12-17: 6 duplicates of slide 2
    new_order = [
        0,   # 1: Title
        1,   # 2: Aktualnist
        2,   # 3: Timeline
        3,   # 4: Problem
        12,  # 5: NEW models
        4,   # 6: Methodology
        13,  # 7: NEW metrics formulas
        14,  # 8: NEW MCDA methodology
        5,   # 9: Architecture
        6,   # 10: Software
        7,   # 11: Dataset
        17,  # 12: NEW Sample queries
        8,   # 13: Results nDCG
        15,  # 14: NEW detailed metrics
        9,   # 15: Analysis (MCDA results)
        16,  # 16: NEW quality vs speed
        10,  # 17: Approbation
        11,  # 18: Summary
    ]
    reorder_slides(prs, new_order)
    print(f"Reordered: {len(prs.slides)} slides")

    # 3. Заповнюємо контент
    fillers = [
        fill_slide_title,
        fill_slide_actuality,
        fill_slide_timeline,
        fill_slide_problem,
        fill_slide_models,
        fill_slide_methodology,
        fill_slide_metrics_formulas,
        fill_slide_mcda,
        fill_slide_architecture,
        fill_slide_software,
        fill_slide_dataset,
        fill_slide_sample_queries,
        fill_slide_results,
        fill_slide_detailed_metrics,
        fill_slide_analysis,
        fill_slide_quality_speed,
        fill_slide_approbation,
        fill_slide_summary,
    ]
    for i, fn in enumerate(fillers):
        slide = prs.slides[i]
        try:
            fn(slide)
            print(f"  Slide {i+1}: filled by {fn.__name__}")
        except Exception as e:
            print(f"  Slide {i+1}: ERROR {fn.__name__}: {e}")
            raise

    # 4. Update slide numbers
    for i, slide in enumerate(prs.slides):
        if i == 0:
            continue  # title slide doesn't show number
        set_slide_number(slide, i + 1)

    # 5. Save
    prs.save(str(PPTX))
    print(f"Saved: {PPTX} ({PPTX.stat().st_size:,} bytes)")
    print(f"Total slides: {len(prs.slides)}")


if __name__ == "__main__":
    main()
